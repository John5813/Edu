"""
Presentation Templates Service
Handles background templates for presentations
"""

import os
import logging
from typing import Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import io

logger = logging.getLogger(__name__)

class PresentationTemplates:
    """Manages presentation background templates"""
    
    def __init__(self):
        self.templates_dir = "webapp/templates"
        self.templates = {
            'classic': {
                'name': 'Klassik',
                'background': None,
                'text_color': RGBColor(0, 0, 0),
                'title_color': RGBColor(0, 51, 102)
            },
            'vintage': {
                'name': 'Vintage',
                'background': 'vintage.jpg',
                'text_color': RGBColor(51, 51, 51),
                'title_color': RGBColor(102, 51, 0)
            },
            'floral': {
                'name': 'Gul naqshli',
                'background': 'floral.jpg',
                'text_color': RGBColor(102, 51, 102),
                'title_color': RGBColor(153, 51, 102)
            },
            'gradient-green': {
                'name': 'Yashil gradient',
                'background': 'gradient-green.jpg',
                'text_color': RGBColor(0, 51, 51),
                'title_color': RGBColor(0, 102, 102)
            },
            'gradient-blue': {
                'name': "Ko'k gradient",
                'background': 'gradient-blue.jpg',
                'text_color': RGBColor(0, 51, 102),
                'title_color': RGBColor(0, 102, 204)
            },
            'gradient-purple': {
                'name': 'Binafsha gradient',
                'background': 'gradient-purple.jpg',
                'text_color': RGBColor(51, 0, 102),
                'title_color': RGBColor(102, 51, 153)
            },
            'poppy': {
                'name': 'Lola gullari',
                'background': 'poppy.jpg',
                'text_color': RGBColor(102, 51, 0),
                'title_color': RGBColor(153, 0, 0)
            },
            'tech-blue': {
                'name': 'Texnologiya',
                'background': 'tech-blue.jpg',
                'text_color': RGBColor(255, 255, 255),
                'title_color': RGBColor(0, 204, 255)
            }
        }
    
    def apply_template_to_slide(self, slide, template_name: str):
        """Apply template background to a slide"""
        try:
            template = self.templates.get(template_name, self.templates['classic'])
            
            # Add background image if specified
            if template['background']:
                bg_path = os.path.join(self.templates_dir, template['background'])
                if os.path.exists(bg_path):
                    # Set slide background
                    self._set_slide_background(slide, bg_path)
                else:
                    logger.warning(f"Background image not found: {bg_path}")
            
            return template
            
        except Exception as e:
            logger.error(f"Error applying template: {e}")
            return self.templates['classic']
    
    def _set_slide_background(self, slide, image_path: str):
        """Set background image for a slide"""
        try:
            # Get slide dimensions
            prs = slide.part.presentation
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Add image as background
            pic = slide.shapes.add_picture(
                image_path,
                0, 0,
                width=slide_width,
                height=slide_height
            )
            
            # Send picture to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
            
        except Exception as e:
            logger.error(f"Error setting slide background: {e}")
    
    def get_template_colors(self, template_name: str) -> Dict:
        """Get color scheme for a template"""
        template = self.templates.get(template_name, self.templates['classic'])
        return {
            'text': template['text_color'],
            'title': template['title_color']
        }
    
    def copy_uploaded_images_to_templates(self):
        """Copy uploaded images from attached_assets to templates folder"""
        try:
            # Map uploaded images to template names
            image_mapping = {
                '62379448511baea4b1ab68cd7a4654c3_1755877745837.jpg': 'vintage.jpg',
                'e2cc2f46ec22c6a2455670d6f4d57cab_1755877800462.jpg': 'floral.jpg',
                '200a00e49abe4c831ac600002b1e1dcd_1755877800513.jpg': 'gradient-green.jpg',
                '184c6f428e31410bd29e2abe04a6582c_1755877800564.jpg': 'gradient-blue.jpg',
                '042e750cbc58c97de7c0a07512475e5d_1755877800618.jpg': 'gradient-purple.jpg',
                'e7d2d0af98b5e8a58374d63335615b86_1755877800676.jpg': 'poppy.jpg',
                '1663850540_g-13_1755879411607.jpg': 'tech-blue.jpg'
            }
            
            os.makedirs(self.templates_dir, exist_ok=True)
            
            for source_name, target_name in image_mapping.items():
                source_path = f"attached_assets/{source_name}"
                target_path = os.path.join(self.templates_dir, target_name)
                
                if os.path.exists(source_path):
                    import shutil
                    shutil.copy2(source_path, target_path)
                    logger.info(f"Copied template: {source_name} -> {target_name}")
                else:
                    logger.warning(f"Source image not found: {source_path}")
                    
        except Exception as e:
            logger.error(f"Error copying template images: {e}")

# Initialize and copy templates on module load
template_service = PresentationTemplates()
template_service.copy_uploaded_images_to_templates()