"""
Template Service for Presentation Backgrounds
Manages background template selection and application
"""

import os
import logging
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

class TemplateService:
    """Manages presentation background templates"""
    
    def __init__(self):
        self.templates = {
            'template_1': {
                'name': {'uz': 'Ta\'lim Elementlari', 'ru': 'Образовательные Элементы', 'en': 'Educational Elements'},
                'file': 'IMG_20250823_092830_649_1764043160410.jpg',
                'colors': {'title': RGBColor(0, 102, 204), 'text': RGBColor(51, 51, 51)}
            },
            'template_2': {
                'name': {'uz': 'Ko\'k Geometrik', 'ru': 'Синий Геометрический', 'en': 'Blue Geometric'},
                'file': 'IMG_20250823_092829_196_1764043160430.jpg',
                'colors': {'title': RGBColor(0, 102, 255), 'text': RGBColor(0, 51, 153)}
            },
            'template_3': {
                'name': {'uz': 'Ko\'k Naqsh', 'ru': 'Синий Узор', 'en': 'Blue Pattern'},
                'file': 'IMG_20250823_092826_569_1764043160448.jpg',
                'colors': {'title': RGBColor(0, 153, 255), 'text': RGBColor(0, 102, 153)}
            },
            'template_4': {
                'name': {'uz': 'Vintage Gul', 'ru': 'Винтажные Цветы', 'en': 'Vintage Flowers'},
                'file': 'IMG_20250823_092824_903_1764043160465.jpg',
                'colors': {'title': RGBColor(153, 0, 51), 'text': RGBColor(102, 51, 51)}
            },
            'template_5': {
                'name': {'uz': 'Pushti Gul', 'ru': 'Розовые Цветы', 'en': 'Pink Flowers'},
                'file': 'IMG_20250823_092822_574_1764043160484.jpg',
                'colors': {'title': RGBColor(153, 51, 153), 'text': RGBColor(102, 51, 102)}
            },
            'template_6': {
                'name': {'uz': 'Rangli Olti Burchak', 'ru': 'Цветной Шестиугольник', 'en': 'Colorful Hexagon'},
                'file': 'IMG_20250823_092819_616_1764043160501.jpg',
                'colors': {'title': RGBColor(255, 102, 0), 'text': RGBColor(51, 51, 51)}
            },
            'template_7': {
                'name': {'uz': 'Minimalist', 'ru': 'Минималист', 'en': 'Minimalist'},
                'file': 'IMG_20250823_092810_644_1764043160522.jpg',
                'colors': {'title': RGBColor(0, 102, 153), 'text': RGBColor(51, 51, 51)}
            },
            'template_8': {
                'name': {'uz': 'Ko\'k To\'lqin', 'ru': 'Синяя Волна', 'en': 'Blue Wave'},
                'file': 'IMG_20250823_092808_412_1764043160540.jpg',
                'colors': {'title': RGBColor(0, 102, 204), 'text': RGBColor(0, 51, 102)}
            },
            'template_9': {
                'name': {'uz': 'Professional Ko\'k', 'ru': 'Профессиональный Синий', 'en': 'Professional Blue'},
                'file': 'IMG_20250823_092805_902_1764043160558.jpg',
                'colors': {'title': RGBColor(0, 51, 153), 'text': RGBColor(0, 51, 102)}
            },
            'template_10': {
                'name': {'uz': 'Zamonaviy Ko\'k', 'ru': 'Современный Синий', 'en': 'Modern Blue'},
                'file': 'IMG_20250823_092803_858_1764043160582.jpg',
                'colors': {'title': RGBColor(0, 102, 255), 'text': RGBColor(0, 51, 153)}
            },
            'template_11': {
                'name': {'uz': 'Piksel Ko\'k', 'ru': 'Пиксельный Синий', 'en': 'Pixel Blue'},
                'file': 'IMG_20250823_092801_417_1764043160607.jpg',
                'colors': {'title': RGBColor(51, 153, 255), 'text': RGBColor(0, 102, 204)}
            },
            'template_12': {
                'name': {'uz': 'Klassik Vintage', 'ru': 'Классический Винтаж', 'en': 'Classic Vintage'},
                'file': 'IMG_20250823_092757_637_1764043160632.jpg',
                'colors': {'title': RGBColor(102, 51, 0), 'text': RGBColor(51, 51, 51)}
            },
            'template_13': {
                'name': {'uz': 'Yashil-Sariq', 'ru': 'Зелено-Желтый', 'en': 'Green-Yellow'},
                'file': 'IMG_20250823_092755_564_1764043160658.jpg',
                'colors': {'title': RGBColor(0, 102, 102), 'text': RGBColor(0, 51, 51)}
            },
            'template_14': {
                'name': {'uz': 'Ish Stoli', 'ru': 'Рабочий Стол', 'en': 'Desktop'},
                'file': 'IMG_20250823_092753_435_1764043160682.jpg',
                'colors': {'title': RGBColor(51, 51, 51), 'text': RGBColor(102, 102, 102)}
            },
            'template_15': {
                'name': {'uz': 'Ta\'lim Buyumlari', 'ru': 'Учебные Предметы', 'en': 'School Supplies'},
                'file': 'IMG_20250823_092752_204_1764043160709.jpg',
                'colors': {'title': RGBColor(0, 153, 255), 'text': RGBColor(0, 102, 153)}
            },
            'template_16': {
                'name': {'uz': 'Texnologiya', 'ru': 'Технология', 'en': 'Technology'},
                'file': 'IMG_20250823_092749_483_1764043160737.jpg',
                'colors': {'title': RGBColor(0, 102, 255), 'text': RGBColor(0, 51, 153)}
            },
            'template_17': {
                'name': {'uz': 'Yashil Gradient', 'ru': 'Зеленый Градиент', 'en': 'Green Gradient'},
                'file': 'IMG_20250823_092748_082_1764043160767.jpg',
                'colors': {'title': RGBColor(0, 153, 102), 'text': RGBColor(0, 102, 51)}
            },
            'template_18': {
                'name': {'uz': 'Bahor Gullari', 'ru': 'Весенние Цветы', 'en': 'Spring Flowers'},
                'file': 'IMG_20250823_092737_363_1764043160796.jpg',
                'colors': {'title': RGBColor(255, 51, 102), 'text': RGBColor(102, 51, 51)}
            },
            'template_19': {
                'name': {'uz': 'Ko\'k Burish', 'ru': 'Синий Угол', 'en': 'Blue Corner'},
                'file': 'IMG_20250823_092717_373_1764043160824.jpg',
                'colors': {'title': RGBColor(0, 102, 255), 'text': RGBColor(0, 51, 153)}
            },
            'template_20': {
                'name': {'uz': 'Ta\'lim Dizayn', 'ru': 'Образовательный Дизайн', 'en': 'Educational Design'},
                'file': 'word-image-3281-8-13_1764043160854.png',
                'colors': {'title': RGBColor(0, 51, 102), 'text': RGBColor(51, 51, 51)}
            }
        }
    
    def get_template_groups(self) -> List[List[Dict]]:
        """Get templates grouped by 5"""
        templates_list = list(self.templates.items())
        groups = []
        
        for i in range(0, len(templates_list), 5):
            group = []
            for j in range(i, min(i + 5, len(templates_list))):
                template_id, template_data = templates_list[j]
                group.append({
                    'id': template_id,
                    'name': template_data['name'],
                    'file': template_data['file']
                })
            groups.append(group)
        
        return groups
    
    def apply_template_to_slide(self, slide, template_id: str):
        """Apply template background to a slide"""
        try:
            if template_id not in self.templates:
                template_id = 'template_20'  # Default
                
            template = self.templates[template_id]
            
            # Add background image if specified
            if template['file']:
                bg_path = os.path.join('attached_assets', template['file'])
                if os.path.exists(bg_path):
                    self._set_slide_background(slide, bg_path)
                else:
                    logger.warning(f"Background image not found: {bg_path}")
            
            return template
            
        except Exception as e:
            logger.error(f"Error applying template: {e}")
            return self.templates['template_20']  # Default
    
    def _set_slide_background(self, slide, image_path: str):
        """Set background image for a slide"""
        try:
            from pptx.util import Inches
            
            # Use standard slide dimensions (16:9)
            slide_width = Inches(13.33)
            slide_height = Inches(7.5)
            
            # Add background image at position 0,0 filling entire slide
            pic = slide.shapes.add_picture(
                image_path,
                0, 0,
                width=slide_width,
                height=slide_height
            )
            
            # Move background to the very back (behind all content)
            # Get all shapes in slide
            shapes_tree = slide.shapes._spTree
            
            # Remove the picture element
            shapes_tree.remove(pic._element)
            
            # Insert it at the beginning (after layout elements)
            # Index 0 and 1 are usually slide layout elements, so insert at 2
            if len(shapes_tree) >= 2:
                shapes_tree.insert(2, pic._element)  
            else:
                shapes_tree.insert(0, pic._element)
            
            logger.info(f"Successfully applied background image: {image_path}")
            
        except Exception as e:
            logger.error(f"Error setting slide background: {e}")
            # Alternative approach if first method fails
            try:
                from pptx.util import Inches
                slide.shapes.add_picture(
                    image_path,
                    0, 0, 
                    Inches(13.33), Inches(7.5)
                )
                logger.info(f"Applied background with alternative method: {image_path}")
            except Exception as e2:
                logger.error(f"Alternative background method also failed: {e2}")
    
    def get_template_colors(self, template_id: str) -> Dict:
        """Get color scheme for a template"""
        template = self.templates.get(template_id, self.templates['template_20'])
        return template['colors']
    
    def get_template_name(self, template_id: str, language: str = 'uz') -> str:
        """Get template name in specified language"""
        template = self.templates.get(template_id, self.templates['template_20'])
        name_dict = template['name']
        
        # Return name based on language, fallback to uzbek
        if isinstance(name_dict, dict):
            return name_dict.get(language, name_dict.get('uz', 'Standart'))
        else:
            # For backward compatibility if name is still string
            return name_dict