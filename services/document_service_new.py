import os
import logging
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from typing import Dict, List, Optional
import asyncio
from bot.services.pexels import PexelsService
from services.ai_service_new import AIService

logger = logging.getLogger(__name__)

class DocumentService:
    def __init__(self, documents_dir: str = "generated_documents"):
        self.documents_dir = documents_dir
        self.ai_service = AIService()
        
        # Ensure directories exist
        os.makedirs(self.documents_dir, exist_ok=True)
        os.makedirs("temp", exist_ok=True)

    async def create_presentation_with_template_background(self, topic: str, content: Dict, author_name: str, template_id: str, template_service) -> str:
        """Create presentation with template background applied"""
        try:
            # First create normal presentation
            temp_file = await self.create_new_presentation_system(topic, content, author_name)
            
            # Now apply template backgrounds to all slides
            from pptx import Presentation
            prs = Presentation(temp_file)
            
            logger.info(f"Applying template {template_id} to {len(prs.slides)} slides")
            
            # Apply template background to each slide
            for slide in prs.slides:
                template_service.apply_template_to_slide(slide, template_id)
                logger.info(f"Applied template background to slide")
            
            # Save with template name
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            template_name = template_service.templates.get(template_id, {}).get('name', 'standart')
            filename = f"template_{template_name}_{timestamp}.pptx"
            file_path = os.path.join(self.documents_dir, filename)
            
            prs.save(file_path)
            
            # Remove temporary file
            if os.path.exists(temp_file):
                os.remove(temp_file)
            
            logger.info(f"Template presentation saved: {file_path}")
            return file_path
            
        except Exception as e:
            logger.error(f"Error creating template presentation: {e}")
            # Fallback to regular presentation
            return await self.create_new_presentation_system(topic, content, author_name)

    async def _apply_template_background(self, slide, template_id: str, template_service):
        """Apply template background image to slide"""
        try:
            template_info = template_service.templates.get(template_id, {})
            template_file = template_info.get('file')
            
            if template_file:
                image_path = f"attached_assets/{template_file}"
                if os.path.exists(image_path):
                    # Add background image to fill entire slide
                    slide.shapes.add_picture(
                        image_path,
                        PptxInches(0), PptxInches(0),
                        PptxInches(13.33), PptxInches(7.5)
                    )
                    # Move background to back of slide
                    slide.shapes[0]._element.getparent().insert(0, slide.shapes[0]._element)
                    logger.info(f"Applied template background: {template_file}")
                    
        except Exception as e:
            logger.warning(f"Could not apply template background: {e}")

    async def _create_title_slide_with_template(self, prs, topic: str, author_name: str, template_service, template_id: str):
        """Create title slide with template background and colors"""
        from pptx.util import Inches as PptxInches, Pt as PptxPt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Apply template background
        await self._apply_template_background(slide, template_id, template_service)
        
        # Get template colors
        colors = template_service.get_template_colors(template_id)
        
        # Add title with template colors
        title_box = slide.shapes.add_textbox(
            PptxInches(1), PptxInches(2.5),
            PptxInches(11.33), PptxInches(2)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = topic
        title_para.font.size = PptxPt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = colors.get('title', RGBColor(0, 51, 102))
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add author
        author_box = slide.shapes.add_textbox(
            PptxInches(1), PptxInches(5),
            PptxInches(11.33), PptxInches(1)
        )
        author_frame = author_box.text_frame
        author_para = author_frame.paragraphs[0]
        author_para.text = f"Muallif: {author_name}"
        author_para.font.size = PptxPt(24)
        author_para.font.color.rgb = colors.get('text', RGBColor(51, 51, 51))
        author_para.alignment = PP_ALIGN.CENTER

    async def _create_content_slide_with_template(self, prs, slide_data: Dict, layout_type: str, slide_num: int, images: Dict, template_service, template_id: str):
        """Create content slide with template background"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Apply template background
        await self._apply_template_background(slide, template_id, template_service)
        
        # Create content based on layout type using existing methods but with template colors
        if layout_type == "bullet_points":
            await self._create_template_bullet_slide(slide, slide_data, template_service, template_id)
        elif layout_type == "text_with_image":
            await self._create_template_text_image_slide(slide, slide_data, slide_num, images, template_service, template_id)
        else:  # three_column
            await self._create_template_three_column_slide(slide, slide_data, template_service, template_id)

    async def _create_template_bullet_slide(self, slide, slide_data: Dict, template_service, template_id: str):
        """Create bullet points slide with template colors"""
        from pptx.util import Inches as PptxInches, Pt as PptxPt
        from pptx.enum.text import PP_ALIGN
        
        colors = template_service.get_template_colors(template_id)
        
        # Title
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.5),
            PptxInches(12), PptxInches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get('title', 'Asosiy Nuqtalar')
        title_para.font.size = PptxPt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = colors.get('title', RGBColor(0, 51, 102))
        title_para.alignment = PP_ALIGN.CENTER
        
        # Content with bullet points using existing logic
        content_box = slide.shapes.add_textbox(
            PptxInches(1), PptxInches(2),
            PptxInches(11.33), PptxInches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Parse content into bullet points using existing method
        content_text = slide_data.get('content', '')
        bullet_points = self._parse_bullet_points(content_text)
        
        for point in bullet_points[:5]:  # Max 5 points
            p = content_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = PptxPt(16)  # Professional small font
            p.font.color.rgb = colors.get('text', RGBColor(51, 51, 51))
            p.level = 0

    async def _create_template_text_image_slide(self, slide, slide_data: Dict, slide_num: int, images: Dict, template_service, template_id: str):
        """Create text+image slide with template colors"""
        from pptx.util import Inches as PptxInches, Pt as PptxPt
        from pptx.enum.text import PP_ALIGN
        
        colors = template_service.get_template_colors(template_id)
        
        # Title
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.5),
            PptxInches(12), PptxInches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get('title', 'Matn va Rasm')
        title_para.font.size = PptxPt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = colors.get('title', RGBColor(0, 51, 102))
        title_para.alignment = PP_ALIGN.CENTER
        
        # Left side: Text (45% width)
        text_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(2),
            PptxInches(5.5), PptxInches(4.5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.text = slide_data.get('content', 'Mazmun mavjud emas')
        text_para.font.size = PptxPt(18)
        text_para.font.color.rgb = colors.get('text', RGBColor(51, 51, 51))
        text_para.alignment = PP_ALIGN.LEFT
        
        # Right side: DALL-E image (55% width)
        if slide_num in images:
            image_path = images[slide_num]
            logger.info(f"Adding DALL-E image for slide {slide_num}: {image_path}")
            if image_path and os.path.exists(image_path):
                try:
                    slide.shapes.add_picture(
                        image_path,
                        PptxInches(6.2), PptxInches(2),    # Right side position
                        PptxInches(6.8), PptxInches(4.5)   # Full right side coverage
                    )
                    logger.info(f"Successfully added DALL-E image to slide {slide_num}")
                except Exception as e:
                    logger.error(f"Error adding DALL-E image to slide {slide_num}: {e}")

    async def _create_template_three_column_slide(self, slide, slide_data: Dict, template_service, template_id: str):
        """Create three column slide with template colors"""
        from pptx.util import Inches as PptxInches, Pt as PptxPt
        from pptx.enum.text import PP_ALIGN
        
        colors = template_service.get_template_colors(template_id)
        
        # Title
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.5),
            PptxInches(12), PptxInches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get('title', 'Uch Ustunli Slayd')
        title_para.font.size = PptxPt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = colors.get('title', RGBColor(0, 51, 102))
        title_para.alignment = PP_ALIGN.CENTER
        
        # Parse content into 3 columns using existing method
        content_text = slide_data.get('content', '')
        columns = self._parse_three_columns_smart(content_text, slide_data.get('title', ''))
        
        # Create 3 columns with template colors
        for i, column in enumerate(columns[:3]):
            x_pos = PptxInches(0.5 + i * 4.2)
            
            col_box = slide.shapes.add_textbox(
                x_pos, PptxInches(2), 
                PptxInches(3.8), PptxInches(4.5)
            )
            col_frame = col_box.text_frame
            col_frame.word_wrap = True
            
            # Column title with template colors
            col_para = col_frame.paragraphs[0]
            col_para.text = column.get('title', f'Ustun {i+1}')
            col_para.font.size = PptxPt(16)
            col_para.font.bold = True
            col_para.font.color.rgb = colors.get('title', RGBColor(0, 51, 102))
            col_para.alignment = PP_ALIGN.CENTER
            
            # Column content with template colors
            column_text = column.get('text', 'Ma\'lumot yo\'q')
            p = col_frame.add_paragraph()
            p.text = str(column_text).strip()
            p.font.size = PptxPt(12)
            p.font.color.rgb = colors.get('text', RGBColor(51, 51, 51))
            p.alignment = PP_ALIGN.LEFT

    async def create_new_presentation_system(self, topic: str, content: Dict, author_name: str) -> str:
        """Create presentation with new 3-template rotating system and DALL-E images"""
        try:
            # Validate content
            if not content or 'slides' not in content:
                logger.error(f"Invalid content structure: {content}")
                raise ValueError("Content must contain 'slides' key")
            
            prs = Presentation()
            
            # Set slide size (16:9)
            prs.slide_width = PptxInches(13.33)
            prs.slide_height = PptxInches(7.5)
            
            slides_data = content.get('slides', [])
            logger.info(f"Creating presentation with {len(slides_data)} slides")
            
            # Generate DALL-E images for text+image slides
            images = await self._generate_dalle_images_for_slides(topic, slides_data)
            
            for idx, slide_data in enumerate(slides_data):
                slide_num = slide_data.get('slide_number', idx + 1)
                layout_type = slide_data.get('layout_type', 'bullet_points')
                
                logger.info(f"Creating slide {slide_num} with layout: {layout_type}")
                
                if slide_num == 1 or layout_type == "title":
                    await self._create_title_slide(prs, topic, author_name)
                else:
                    await self._create_new_content_slide(prs, slide_data, layout_type, slide_num, images)
            
            # Save presentation
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"new_presentation_{timestamp}.pptx"
            file_path = os.path.join(self.documents_dir, filename)
            
            prs.save(file_path)
            logger.info(f"New presentation system saved: {file_path}")
            
            return file_path
            
        except Exception as e:
            logger.error(f"Error creating new presentation: {e}")
            raise

    async def _generate_dalle_images_for_slides(self, topic: str, slides_data: List[Dict]) -> Dict[int, str]:
        """Generate DALL-E images for text+image layout slides"""
        images_dict = {}
        
        try:
            for slide_data in slides_data:
                slide_num = slide_data.get('slide_number', 0)
                layout_type = slide_data.get('layout_type', '')
                
                if layout_type == "text_with_image":
                    slide_title = slide_data.get('title', '')
                    slide_content = slide_data.get('content', '')
                    
                    # Generate DALL-E image
                    image_url = await self.ai_service.generate_dalle_image(
                        slide_content, slide_title
                    )
                    
                    if image_url:
                        # Download image
                        filename = f"dalle_slide_{slide_num}.png"
                        image_path = await self.ai_service.download_image(image_url, filename)
                        
                        if image_path:
                            images_dict[slide_num] = image_path
                            logger.info(f"Generated DALL-E image for slide {slide_num}: {slide_title}")
                    
                    # Delay between image generations
                    await asyncio.sleep(1.0)
            
            return images_dict
            
        except Exception as e:
            logger.error(f"Error generating DALL-E images: {e}")
            return {}

    async def _create_new_content_slide(self, prs, slide_data: Dict, layout_type: str, slide_num: int, images: Dict):
        """Create content slide with new system layouts"""
        logger.info(f"Creating slide {slide_num} with layout '{layout_type}', title: '{slide_data.get('title', 'NO TITLE')}', content length: {len(slide_data.get('content', ''))}")
        
        if layout_type == "bullet_points":
            await self._create_new_bullet_points_slide(prs, slide_data)
        elif layout_type == "text_with_image":
            await self._create_new_text_with_image_slide(prs, slide_data, slide_num, images)
        elif layout_type == "three_column":
            await self._create_new_three_column_slide(prs, slide_data)
        else:
            logger.error(f"Unknown layout type: {layout_type}, using bullet_points fallback")
            await self._create_new_bullet_points_slide(prs, slide_data)

    async def _create_new_bullet_points_slide(self, prs, slide_data: Dict):
        """Create LAYOUT 1: 5 bullet points with 30+ words each (2,5,8,11...)"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        # Title
        if title:
            title.text = slide_data.get('title', 'Bullet Points Slayd')
            title.text_frame.paragraphs[0].font.size = PptxPt(28)  # Medium title
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Content - parse into bullet points
        if content_placeholder:
            content_text = slide_data.get('content', '')
            content_frame = content_placeholder.text_frame
            content_frame.clear()
            
            # Split content into bullet points
            bullet_points = self._parse_bullet_points(content_text)
            
            for i, point in enumerate(bullet_points[:5]):  # Max 5 points
                if i == 0:
                    # First paragraph
                    p = content_frame.paragraphs[0]
                else:
                    # Additional paragraphs
                    p = content_frame.add_paragraph()
                
                p.text = f"• {point.strip()}"
                p.font.size = PptxPt(14)  # Small font for detailed content
                p.alignment = PP_ALIGN.LEFT
                p.level = 0

    async def _create_new_text_with_image_slide(self, prs, slide_data: Dict, slide_num: int, images: Dict):
        """Create LAYOUT 2: Text + DALL-E image (50/50 split) (3,6,9,12...)"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Title background for better visibility
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            PptxInches(0.3), PptxInches(0.3),
            PptxInches(12.4), PptxInches(1.4)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
        title_bg.fill.transparency = 0.3  # 30% transparency
        title_bg.line.fill.background()  # No border
        
        # Title
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.5),
            PptxInches(12), PptxInches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get('title', 'Matn + Rasm Slayd')
        title_para.font.size = PptxPt(32)  # Larger title
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)  # White text
        title_para.alignment = PP_ALIGN.CENTER

        # Add semi-transparent background for text readability
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        
        text_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            PptxInches(0.3), PptxInches(1.8),
            PptxInches(5.9), PptxInches(4.9)
        )
        text_bg.fill.solid()
        text_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
        text_bg.fill.transparency = 0.2  # 20% transparency
        text_bg.line.fill.background()  # No border
        
        # Left side: Continuous text (45% width)
        text_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(2),
            PptxInches(5.5), PptxInches(4.5)  # 45% width
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.text = slide_data.get('content', 'Mazmun mavjud emas')
        text_para.font.size = PptxPt(18)  # Large font for main content
        text_para.font.bold = True  # Bold for better visibility
        text_para.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        text_para.alignment = PP_ALIGN.LEFT  # Left alignment

        # Right side: DALL-E image (55% width - butun o'ng tomonni qoplash)
        if slide_num in images:
            image_path = images[slide_num]
            logger.info(f"Adding DALL-E image for slide {slide_num}: {image_path}")
            if image_path and os.path.exists(image_path):
                try:
                    slide.shapes.add_picture(
                        image_path,
                        PptxInches(6.2), PptxInches(2),    # Right side position
                        PptxInches(6.8), PptxInches(4.5)   # O'ng tomonni butunlay qoplash
                    )
                    logger.info(f"Successfully added DALL-E image to slide {slide_num}")
                except Exception as e:
                    logger.error(f"Error adding DALL-E image to slide {slide_num}: {e}")
            else:
                logger.warning(f"DALL-E image not found for slide {slide_num}: {image_path}")
        else:
            logger.info(f"No DALL-E image available for slide {slide_num}")

    async def _create_new_three_column_slide(self, prs, slide_data: Dict):
        """Create LAYOUT 3: Smart 3-column layout with logical headers (4,7,10,13...)"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Title background for better visibility
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            PptxInches(0.3), PptxInches(0.3),
            PptxInches(12.4), PptxInches(1.4)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background
        title_bg.fill.transparency = 0.3  # 30% transparency
        title_bg.line.fill.background()  # No border
        
        # Title
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.5),
            PptxInches(12), PptxInches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data.get('title', 'Uch Ustunli Slayd')
        title_para.font.size = PptxPt(32)  # Larger title
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)  # White text
        title_para.alignment = PP_ALIGN.CENTER

        # Parse content into 3 logical columns
        content_text = slide_data.get('content', '')
        # Handle case where content might be a dict or other type
        if not isinstance(content_text, str):
            content_text = str(content_text) if content_text else ''
        
        columns = self._parse_three_columns_smart(content_text, slide_data.get('title', ''))
        
        # Create 3 columns
        column_width = PptxInches(3.8)
        column_height = PptxInches(4.5)
        start_x = PptxInches(0.5)
        start_y = PptxInches(2)

        for i, column in enumerate(columns[:3]):
            # Calculate position
            x_pos = start_x + i * PptxInches(4.2)
            
            # Add column textbox
            col_box = slide.shapes.add_textbox(x_pos, start_y, column_width, column_height)
            col_frame = col_box.text_frame
            col_frame.word_wrap = True
            
            # Column title (logical header)
            col_para = col_frame.paragraphs[0]
            col_para.text = column.get('title', f'Ustun {i+1}')
            col_para.font.size = PptxPt(16)  # Medium font for headers
            col_para.font.bold = True
            col_para.font.name = 'Times New Roman'  # Times New Roman shrift
            col_para.font.color.rgb = RGBColor(255, 255, 255)  # White text for visibility
            col_para.alignment = PP_ALIGN.CENTER
            
            # Column continuous text (40 words without bullets or numbers)
            column_text = column.get('text', column.get('points', ['Ma\'lumot yo\'q'])[0] if column.get('points') else 'Ma\'lumot yo\'q')
            
            # Add continuous text paragraph
            p = col_frame.add_paragraph()
            p.text = str(column_text).strip()
            p.font.size = PptxPt(12)  # Small font for details
            p.font.name = 'Times New Roman'  # Times New Roman shrift
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black text
            p.alignment = PP_ALIGN.LEFT  # Left alignment
            p.level = 0

    def _parse_bullet_points(self, content_text: str) -> List[str]:
        """Parse content into bullet points (aim for 5 points with 30+ words each)"""
        # Ensure content_text is string
        if not isinstance(content_text, str):
            content_text = str(content_text) if content_text else ''
            
        if not content_text or content_text.strip() == '':
            return ["Ma'lumot mavjud emas"] * 5
        
        # Try to split by existing bullet points or numbers
        points = []
        
        # Check for existing bullet points
        if '•' in content_text:
            points = [p.strip() for p in content_text.split('•') if p.strip()]
        elif '\n-' in content_text:
            points = [p.strip() for p in content_text.split('\n-') if p.strip()]
        elif '\n' in content_text and len(content_text.split('\n')) >= 3:
            points = [p.strip() for p in content_text.split('\n') if p.strip()]
        else:
            # Split by sentences
            sentences = [s.strip() for s in content_text.split('.') if s.strip()]
            if len(sentences) >= 5:
                points = sentences[:5]
            else:
                # Split longer content into chunks
                words = content_text.split()
                chunk_size = max(30, len(words) // 5)  # At least 30 words per point
                points = [' '.join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]
        
        # Ensure we have exactly 5 points
        while len(points) < 5:
            points.append(f"Qo'shimcha ma'lumot {len(points) + 1}")
        
        return points[:5]

    def _parse_three_columns_smart(self, content_text: str, slide_title: str) -> List[Dict]:
        """Parse content into 3 logical columns with 40-word continuous text per column"""
        # Ensure content_text is string
        if not isinstance(content_text, str):
            content_text = str(content_text) if content_text else ''
            
        # Check if content has ||| separator (new AI format)
        if '|||' in content_text:
            columns_data = [col.strip() for col in content_text.split('|||')]
            result = []
            default_titles = ['Asosiy Jihat', 'Qo\'shimcha Ma\'lumot', 'Muhim Nuqta']
            
            # Parse in pairs: title|||content|||title|||content|||title|||content
            for i in range(0, min(len(columns_data), 6), 2):  # Process in pairs
                if i + 1 < len(columns_data):
                    title = columns_data[i].strip()
                    text = columns_data[i + 1].strip()
                else:
                    title = default_titles[i // 2] if i // 2 < len(default_titles) else f'Ustun {i // 2 + 1}'
                    text = columns_data[i].strip() if i < len(columns_data) else f'80 so\'zlik batafsil ma\'lumot kerak.'
                
                # Ensure text is long enough (should be 80+ words)
                if len(text.split()) < 20:  # If less than 20 words, pad it
                    text += f' Ushbu ustun bo\'yicha qo\'shimcha batafsil ma\'lumotlar va tushuntirishlar kiritilishi kerak. Professional akademik uslubda yozilgan to\'liq mazmun bu yerda bo\'lishi lozim.'
                
                result.append({'title': title, 'text': text})
            
            # Fill remaining columns if needed
            while len(result) < 3:
                i = len(result)
                result.append({
                    'title': default_titles[i] if i < len(default_titles) else f'Ustun {i+1}',
                    'text': f'Bu ustun uchun 80 so\'zlik batafsil ma\'lumot kerak. Professional akademik uslubda yozilgan to\'liq tushuntirish va misollar bilan boyitilgan mazmun bu yerda bo\'lishi lozim. Har qanday qo\'shimcha ma\'lumotlar ham qo\'shilishi mumkin.'
                })
            
            return result
            
        if not content_text or content_text.strip() == '':
            return [
                {'title': 'Sabablari', 'text': 'Bu ustun haqida batafsil ma\'lumot mavjud emas. Keyinchalik qo\'shimcha ma\'lumotlar bilan to\'ldiriladi va yangi tushunchalar kiritiladi asosiy sabablari haqida ma\'lumotlar.'},
                {'title': 'Ta\'siri', 'text': 'Ushbu bo\'lim bo\'yicha qo\'shimcha tafsilotlar hali tayyor emas. Vaqt o\'tishi bilan muhim nuqtalar va asosiy ma\'lumotlar qo\'shiladi ta\'sir ko\'rsatuvchi omillar haqida.'},
                {'title': 'Yechimlar', 'text': 'Yakuniy xulosalar va umumlashtiruvchi fikrlar hali shakllantirilmagan. Kelajakda barcha ma\'lumotlar asosida natijalar chiqariladi va tavsiyalar beriladi yechimlar haqida.'}
            ]
        
        # Generate logical headers based on slide title
        headers = self._generate_logical_headers(slide_title)
        
        # Clean content from any bullets or numbers
        clean_content = content_text.replace('•', '').replace('-', '').replace('*', '')
        # Remove numbered lists (1. 2. 3. etc.)
        import re
        clean_content = re.sub(r'\d+\.\s*', '', clean_content)
        
        # Split content into 3 equal parts by word count
        words = clean_content.split()
        total_words = len(words)
        
        if total_words >= 120:  # If we have enough words (3 x 40)
            # Split into 3 equal parts
            words_per_column = total_words // 3
            columns = []
            for i in range(3):
                start_idx = i * words_per_column
                if i == 2:  # Last column gets remaining words
                    end_idx = total_words
                else:
                    end_idx = (i + 1) * words_per_column
                
                column_text = ' '.join(words[start_idx:end_idx])
                columns.append({
                    'title': headers[i],
                    'text': column_text
                })
        else:
            # If not enough content, split sentences
            sentences = [s.strip() for s in clean_content.split('.') if s.strip()]
            if len(sentences) >= 3:
                per_column = max(1, len(sentences) // 3)
                columns = []
                for i in range(3):
                    start_idx = i * per_column
                    end_idx = (i + 1) * per_column if i < 2 else len(sentences)
                    column_sentences = sentences[start_idx:end_idx]
                    column_text = '. '.join(column_sentences) + '.'
                    
                    columns.append({
                        'title': headers[i],
                        'text': column_text
                    })
            else:
                # Fallback: distribute content evenly
                third = len(clean_content) // 3
                columns = [
                    {'title': headers[0], 'text': clean_content[:third]},
                    {'title': headers[1], 'text': clean_content[third:2*third]},
                    {'title': headers[2], 'text': clean_content[2*third:]}
                ]
        
        return columns

    def _generate_logical_headers(self, slide_title: str) -> List[str]:
        """Generate logical headers based on slide title instead of generic 'Qism 1,2,3'"""
        title_lower = slide_title.lower()
        
        # Economics/Business related
        if any(word in title_lower for word in ['iqtisod', 'biznes', 'moliya', 'bozor', 'savdo']):
            return ['Sabablari', 'Ta\'siri', 'Yechimlar']
        
        # Problem/Solution related
        if any(word in title_lower for word in ['muammo', 'masala', 'yechim', 'hal']):
            return ['Muammolar', 'Sabablar', 'Yechimlar']
        
        # Historical/Development related
        if any(word in title_lower for word in ['tarix', 'rivojlanish', 'o\'sish', 'davr']):
            return ['Boshlang\'ich', 'Rivojlanish', 'Natijalar']
        
        # Analysis/Research related
        if any(word in title_lower for word in ['tahlil', 'tadqiqot', 'o\'rganish']):
            return ['Ma\'lumotlar', 'Tahlil', 'Xulosalar']
        
        # Technology related
        if any(word in title_lower for word in ['texnologiya', 'innovatsiya', 'raqamli']):
            return ['Texnologiya', 'Qo\'llanish', 'Kelajak']
        
        # Default logical headers
        return ['Asosiy Jihat', 'Muhim Omil', 'Yakuniy Natija']

    async def _create_title_slide(self, prs, topic: str, author_name: str):
        """Create title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        if title:
            title.text = "Taqdimot"
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title.text_frame.paragraphs[0].font.size = PptxPt(36)

        if subtitle:
            subtitle.text = f"{topic}\n\n\n{author_name or '__________________'}"
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle.text_frame.paragraphs[0].font.size = PptxPt(20)

    async def create_referat(self, topic: str, content: Dict) -> str:
        """Create referat document"""
        try:
            from docx import Document
            from docx.shared import Pt, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()

            # Set document style
            style = doc.styles['Normal']
            font = style.font
            font.name = 'Times New Roman'
            font.size = Pt(12)

            # Title page
            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run("REFERAT")
            title_run.font.size = Pt(16)
            title_run.font.bold = True

            doc.add_paragraph()  # Empty line

            topic_para = doc.add_paragraph()
            topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            topic_run = topic_para.add_run(topic)
            topic_run.font.size = Pt(14)
            topic_run.font.bold = True

            # Add page break
            doc.add_page_break()

            # Table of contents
            toc_para = doc.add_paragraph()
            toc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            toc_run = toc_para.add_run("REJA")
            toc_run.font.size = Pt(14)
            toc_run.font.bold = True

            doc.add_paragraph()  # Empty line

            # Add sections to TOC
            sections = content.get('sections', [])
            for idx, section in enumerate(sections, 1):
                toc_item = doc.add_paragraph()
                toc_item.paragraph_format.first_line_indent = Inches(0.5)
                toc_item.add_run(f"{idx}. {section['title']}")

            # Add page break
            doc.add_page_break()

            # Add sections content
            for idx, section in enumerate(sections, 1):
                # Section title
                section_title = doc.add_paragraph()
                section_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
                section_title_run = section_title.add_run(f"{idx}. {section['title']}")
                section_title_run.font.bold = True
                section_title_run.font.size = Pt(14)

                doc.add_paragraph()  # Empty line

                # Section content
                content_para = doc.add_paragraph(section['content'])
                content_para.paragraph_format.first_line_indent = Inches(0.5)
                content_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

                doc.add_paragraph()  # Empty line

            # References
            if content.get('references'):
                doc.add_page_break()

                ref_title = doc.add_paragraph()
                ref_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
                ref_title_run = ref_title.add_run("FOYDALANILGAN ADABIYOTLAR")
                ref_title_run.font.bold = True
                ref_title_run.font.size = Pt(14)

                doc.add_paragraph()  # Empty line

                for idx, ref in enumerate(content['references'], 1):
                    ref_para = doc.add_paragraph()
                    ref_para.paragraph_format.first_line_indent = Inches(-0.5)
                    ref_para.paragraph_format.left_indent = Inches(0.5)
                    ref_para.add_run(f"{idx}. {ref}")

            # Save document
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"referat_{timestamp}.docx"
            file_path = os.path.join(self.documents_dir, filename)
            
            doc.save(file_path)
            logger.info(f"Referat saved: {file_path}")
            
            return file_path
            
        except Exception as e:
            logger.error(f"Error creating referat: {e}")
            raise

    async def create_independent_work(self, topic: str, content: Dict) -> str:
        """Create independent work document"""
        try:
            from docx import Document
            from docx.shared import Pt, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()

            # Set document style
            style = doc.styles['Normal']
            font = style.font
            font.name = 'Times New Roman'
            font.size = Pt(12)

            # Title page
            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run = title_para.add_run("MUSTAQIL ISH")
            title_run.font.size = Pt(16)
            title_run.font.bold = True

            doc.add_paragraph()  # Empty line

            topic_para = doc.add_paragraph()
            topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            topic_run = topic_para.add_run(topic)
            topic_run.font.size = Pt(14)
            topic_run.font.bold = True

            # Add page break
            doc.add_page_break()

            # Table of contents
            toc_para = doc.add_paragraph()
            toc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            toc_run = toc_para.add_run("REJA")
            toc_run.font.size = Pt(14)
            toc_run.font.bold = True

            doc.add_paragraph()  # Empty line

            # Add sections to TOC
            sections = content.get('sections', [])
            for idx, section in enumerate(sections, 1):
                toc_item = doc.add_paragraph()
                toc_item.paragraph_format.first_line_indent = Inches(0.5)
                toc_item.add_run(f"{idx}. {section['title']}")

            # Add page break
            doc.add_page_break()

            # Add sections content
            for idx, section in enumerate(sections, 1):
                # Section title
                section_title = doc.add_paragraph()
                section_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
                section_title_run = section_title.add_run(f"{idx}. {section['title']}")
                section_title_run.font.bold = True
                section_title_run.font.size = Pt(14)

                doc.add_paragraph()  # Empty line

                # Section content
                content_para = doc.add_paragraph(section['content'])
                content_para.paragraph_format.first_line_indent = Inches(0.5)
                content_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

                doc.add_paragraph()  # Empty line

            # References
            if content.get('references'):
                doc.add_page_break()

                ref_title = doc.add_paragraph()
                ref_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
                ref_title_run = ref_title.add_run("FOYDALANILGAN ADABIYOTLAR")
                ref_title_run.font.bold = True
                ref_title_run.font.size = Pt(14)

                doc.add_paragraph()  # Empty line

                for idx, ref in enumerate(content['references'], 1):
                    ref_para = doc.add_paragraph()
                    ref_para.paragraph_format.first_line_indent = Inches(-0.5)
                    ref_para.paragraph_format.left_indent = Inches(0.5)
                    ref_para.add_run(f"{idx}. {ref}")

            # Save document
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"independent_work_{timestamp}.docx"
            file_path = os.path.join(self.documents_dir, filename)
            
            doc.save(file_path)
            logger.info(f"Independent work saved: {file_path}")
            
            return file_path
            
        except Exception as e:
            logger.error(f"Error creating independent work: {e}")
            raise