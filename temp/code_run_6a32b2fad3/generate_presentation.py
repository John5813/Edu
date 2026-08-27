import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

MAVZU = "Avtomatlashtirilgan qurilish loyihalash"
SLIDE_COUNT = 10
OUTPUT_PATH = os.environ.get("PPTX_OUTPUT_PATH", "presentation.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(45, 55, 72)
GRAY = RGBColor(107, 114, 128)
LIGHT_GRAY = RGBColor(229, 231, 235)
MID_GRAY = RGBColor(156, 163, 175)
GREEN = RGBColor(22, 163, 74)
GREEN_DARK = RGBColor(21, 128, 61)
GREEN_LIGHT = RGBColor(187, 247, 208)
BLUE = RGBColor(37, 99, 235)
BLUE_LIGHT = RGBColor(191, 219, 254)
TEAL = RGBColor(13, 148, 136)
SOFT_BG = RGBColor(248, 250, 252)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=1.2, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.color.rgb = line_color if line_color else fill_color
    shp.line.width = Pt(line_width)
    return shp


def add_line(slide, x1, y1, x2, y2, color=LIGHT_GRAY, width=1.5):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln


def add_textbox(slide, x, y, w, h, text, font_size=20, color=DARK, bold=False,
                font_name="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                margin_left=0.08, margin_right=0.08, margin_top=0.05, margin_bottom=0.05):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(margin_right)
    tf.margin_top = Inches(margin_top)
    tf.margin_bottom = Inches(margin_bottom)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return tx


def add_paragraphs_box(slide, x, y, w, h, paragraphs, font_size=18, color=DARK,
                       bullet=True, level_indent=0.22, spacing=1.12, font_name="Aptos"):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.04)
    for i, item in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        if bullet:
            p.bullet = True
        p.space_after = Pt(5)
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = color
    return tx


def add_title(slide, title, subtitle=None, accent_color=GREEN):
    add_rect(slide, Inches(0.45), Inches(0.42), Inches(0.16), Inches(0.95), accent_color, accent_color, 0.5, radius=True)
    add_textbox(slide, Inches(0.78), Inches(0.35), Inches(8.6), Inches(0.7), title, 26, DARK, True)
    add_line(slide, Inches(0.8), Inches(1.05), Inches(12.5), Inches(1.05), LIGHT_GRAY, 1.2)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.08), Inches(8.5), Inches(0.45), subtitle, 11, GRAY, False)
    add_rect(slide, Inches(11.95), Inches(0.38), Inches(0.8), Inches(0.22), GREEN_LIGHT, GREEN_LIGHT, 0.5, radius=True)
    add_rect(slide, Inches(12.18), Inches(0.68), Inches(0.52), Inches(0.16), BLUE_LIGHT, BLUE_LIGHT, 0.5, radius=True)


def add_footer(slide, idx):
    add_line(slide, Inches(0.55), Inches(7.0), Inches(12.75), Inches(7.0), LIGHT_GRAY, 0.8)
    add_textbox(slide, Inches(0.58), Inches(7.02), Inches(4.0), Inches(0.22), MAVZU, 9, MID_GRAY)
    add_textbox(slide, Inches(12.2), Inches(7.0), Inches(0.4), Inches(0.22), str(idx), 9, MID_GRAY, align=PP_ALIGN.RIGHT)


def stat_card(slide, x, y, w, h, title, value, fill=SOFT_BG, accent=GREEN):
    add_rect(slide, x, y, w, h, fill, LIGHT_GRAY, 1.0, radius=True)
    add_rect(slide, x + Inches(0.12), y + Inches(0.15), Inches(0.12), h - Inches(0.3), accent, accent, 0.5, radius=True)
    add_textbox(slide, x + Inches(0.34), y + Inches(0.18), w - Inches(0.45), Inches(0.28), title, 11, GRAY, False)
    add_textbox(slide, x + Inches(0.34), y + Inches(0.44), w - Inches(0.45), Inches(0.42), value, 20, DARK, True)


def process_node(slide, x, y, w, h, title, body, accent=GREEN):
    add_rect(slide, x, y, w, h, WHITE, LIGHT_GRAY, 1.1, radius=True)
    add_rect(slide, x, y, w, Inches(0.18), accent, accent, 0.5, radius=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.22), w - Inches(0.3), Inches(0.32), title, 16, DARK, True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.56), w - Inches(0.3), h - Inches(0.65), body, 12, GRAY)


def draw_circle(slide, cx, cy, diameter, fill, line=WHITE, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - diameter / 2, cy - diameter / 2, diameter, diameter)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_w)
    return shp


def add_ring_segment(slide, cx, cy, outer_d, width_d, start_deg, end_deg, color, steps=36):
    points_outer = []
    points_inner = []
    for i in range(steps + 1):
        a = math.radians(start_deg + (end_deg - start_deg) * i / steps)
        points_outer.append((cx + math.cos(a) * outer_d / 2, cy + math.sin(a) * outer_d / 2))
    inner_d = outer_d - width_d
    for i in range(steps + 1):
        a = math.radians(end_deg - (end_deg - start_deg) * i / steps)
        points_inner.append((cx + math.cos(a) * inner_d / 2, cy + math.sin(a) * inner_d / 2))
    free = slide.shapes.build_freeform(points_outer[0][0], points_outer[0][1])
    for px, py in points_outer[1:]:
        free.add_line_segments([(px, py)], close=False)
    for px, py in points_inner:
        free.add_line_segments([(px, py)], close=False)
    free_shape = free.convert_to_shape()
    free_shape.fill.solid()
    free_shape.fill.fore_color.rgb = color
    free_shape.line.color.rgb = WHITE
    free_shape.line.width = Pt(1.0)
    return free_shape


def donut_chart_slide(slide, x, y, size, data, colors, center_title, center_subtitle):
    total = sum(v for _, v in data)
    start = -90
    ring_width = size * 0.26
    draw_circle(slide, x + size / 2, y + size / 2, size, LIGHT_GRAY, WHITE, 1.0)
    current = start
    for i, (label, value) in enumerate(data):
        angle = 360 * value / total
        add_ring_segment(slide, x + size / 2, y + size / 2, size, ring_width, current, current + angle, colors[i % len(colors)], 42)
        current += angle
    draw_circle(slide, x + size / 2, y + size / 2, size - ring_width, WHITE, WHITE, 1.0)
    add_textbox(slide, x + Inches(0.65), y + Inches(1.28), Inches(1.8), Inches(0.35), center_title, 18, DARK, True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.55), y + Inches(1.63), Inches(2.0), Inches(0.3), center_subtitle, 11, GRAY, False, align=PP_ALIGN.CENTER)

    lx = x + size + Inches(0.45)
    ly = y + Inches(0.15)
    for i, (label, value) in enumerate(data):
        add_rect(slide, lx, ly + Inches(i * 0.72), Inches(0.18), Inches(0.18), colors[i % len(colors)], colors[i % len(colors)], 0.5, radius=True)
        add_textbox(slide, lx + Inches(0.26), ly + Inches(i * 0.72) - Inches(0.02), Inches(3.1), Inches(0.24),
                    f"{label} — {value}%", 12, DARK, True)
        desc = {
            "BIM modellashtirish": "Loyiha elementlarini raqamli modelga birlashtirish, to‘qnashuvlarni erta topish va ma’lumotlar bazasini shakllantirish.",
            "Parametrik chizmalar": "Bir o‘zgaruvchi almashganda kesim, fasad, reja va spetsifikatsiyalarni avtomatik yangilash imkonini beradi.",
            "Analitik hisoblash": "Konstruksiya, energiya, yuklama va muddat bo‘yicha dastlabki tahlillarni tezlashtiradi.",
            "Nazorat va integratsiya": "Versiyalarni boshqarish, hujjat aylanishi va pudratchilar bilan sinxron ishlashni ta’minlaydi."
        }.get(label, "")
        add_textbox(slide, lx + Inches(0.26), ly + Inches(i * 0.72) + Inches(0.2), Inches(4.0), Inches(0.42), desc, 10, GRAY)


def slide_1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)

    add_rect(slide, Inches(0.55), Inches(0.55), Inches(12.2), Inches(6.35), WHITE, LIGHT_GRAY, 1.2, radius=True)
    add_rect(slide, Inches(0.75), Inches(0.78), Inches(0.18), Inches(5.7), GREEN, GREEN, 0.5, radius=True)
    add_rect(slide, Inches(11.95), Inches(0.8), Inches(0.35), Inches(0.35), BLUE, BLUE, 0.5, radius=True)
    add_rect(slide, Inches(12.35), Inches(1.22), Inches(0.2), Inches(0.2), GREEN_LIGHT, GREEN_LIGHT, 0.5, radius=True)

    add_textbox(slide, Inches(1.18), Inches(1.1), Inches(6.9), Inches(1.4), MAVZU, 28, DARK, True)
    add_textbox(
        slide, Inches(1.2), Inches(2.35), Inches(6.4), Inches(1.25),
        "Qurilish loyihalarini tezroq, aniqroq va boshqariladigan tarzda ishlab chiqish uchun raqamli modellashtirish, parametrik yondashuv va avtomatik nazorat vositalari.",
        16, GRAY
    )

    right = add_rect(slide, Inches(8.3), Inches(1.12), Inches(3.8), Inches(4.85), SOFT_BG, LIGHT_GRAY, 1.0, radius=True)
    add_rect(slide, Inches(8.55), Inches(1.4), Inches(1.05), Inches(1.05), GREEN_LIGHT, GREEN_LIGHT, 0.5, radius=True)
    add_rect(slide, Inches(9.72), Inches(1.4), Inches(1.95), Inches(0.3), BLUE_LIGHT, BLUE_LIGHT, 0.5, radius=True)
    add_rect(slide, Inches(9.72), Inches(1.84), Inches(1.55), Inches(0.24), LIGHT_GRAY, LIGHT_GRAY, 0.5, radius=True)

    for i, h in enumerate([1.0, 1.5, 2.2, 1.7]):
        add_rect(slide, Inches(8.75 + i * 0.72), Inches(4.95 - h), Inches(0.45), Inches(h), [GREEN, TEAL, BLUE, GREEN_DARK][i], [GREEN, TEAL, BLUE, GREEN_DARK][i], 0.5, radius=True)
    add_line(slide, Inches(8.62), Inches(5.02), Inches(11.8), Inches(5.02), MID_GRAY, 1.0)
    add_line(slide, Inches(8.62), Inches(5.02), Inches(8.62), Inches(2.15), MID_GRAY, 1.0)

    add_textbox(slide, Inches(8.55), Inches(5.25), Inches(3.0), Inches(0.5), "Raqamli boshqaruv, BIM, tahlil va muvofiqlashtirish", 11, GRAY)

    stat_card(slide, Inches(1.22), Inches(4.35), Inches(2.2), Inches(1.15), "Asosiy yondashuv", "BIM + CAD")
    stat_card(slide, Inches(3.58), Inches(4.35), Inches(2.2), Inches(1.15), "Natija", "Kam xatolik")
    stat_card(slide, Inches(5.94), Inches(4.35), Inches(2.2), Inches(1.15), "Ta’sir", "Tez qaror")

    add_textbox(slide, Inches(1.22), Inches(5.78), Inches(5.7), Inches(0.48), "Taqdimot tili: o‘zbek", 12, GRAY)
    add_footer(slide, 1)


def slide_2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "1. Sohadagi zarurat va muammo", "Nima uchun avtomatlashtirilgan loyihalash zamonaviy qurilishda strategik ahamiyatga ega?")
    add_rect(slide, Inches(0.68), Inches(1.48), Inches(12.0), Inches(5.15), WHITE, LIGHT_GRAY, 1.0, radius=True)

    add_rect(slide, Inches(0.9), Inches(1.78), Inches(3.8), Inches(4.55), SOFT_BG, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(1.12), Inches(2.0), Inches(3.25), Inches(0.35), "An’anaviy yondashuvdagi cheklovlar", 17, DARK, True)
    add_paragraphs_box(
        slide, Inches(1.08), Inches(2.42), Inches(3.4), Inches(3.5),
        [
            "Qurilish chizmalari ko‘pincha turli mutaxassislar tomonidan alohida ishlab chiqiladi, bu esa ma’lumotlar bir-biriga mos kelmasligiga olib keladi.",
            "Loyihadagi bitta o‘zgarish reja, fasad, kesim va spesifikatsiyalarda qo‘lda qayta tahrir talab qiladi; bu vaqtni uzaytiradi va xatolik ehtimolini oshiradi.",
            "Muddat bosimi ostida tayyorlangan hujjatlarda konstruktiv va muhandislik qismlar orasida to‘qnashuvlar kech aniqlanadi.",
            "Buyurtmachi, pudratchi va loyiha guruhi o‘rtasida yagona raqamli muhit bo‘lmasa, qaror qabul qilish sekinlashadi."
        ],
        12, GRAY
    )

    add_rect(slide, Inches(4.95), Inches(1.78), Inches(3.2), Inches(2.05), WHITE, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(5.18), Inches(2.0), Inches(2.7), Inches(0.3), "Avtomatlashtirishning maqsadi", 16, DARK, True)
    add_textbox(
        slide, Inches(5.18), Inches(2.38), Inches(2.7), Inches(1.15),
        "Loyiha ma’lumotlarini yagona modelda saqlash, parametrlar asosida tez yangilash, hisob-kitoblarni avtomatlashtirish va ko‘p ishtirokchili jarayonni tartibga solish.",
        12, GRAY
    )

    add_rect(slide, Inches(4.95), Inches(4.02), Inches(3.2), Inches(2.3), WHITE, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(5.18), Inches(4.25), Inches(2.7), Inches(0.3), "Kutiladigan foyda", 16, DARK, True)
    add_paragraphs_box(
        slide, Inches(5.12), Inches(4.58), Inches(2.8), Inches(1.45),
        [
            "Hujjat tayyorlash tezligi oshadi.",
            "Loyihalararo izchillik yaxshilanadi.",
            "Qayta ishlash va tuzatish xarajati kamayadi.",
            "Nazorat va tahlil sifati oshadi."
        ],
        11, GRAY
    )

    add_rect(slide, Inches(8.42), Inches(1.78), Inches(4.0), Inches(4.55), SOFT_BG, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(8.68), Inches(2.0), Inches(3.3), Inches(0.35), "Muhim biznes natijalari", 17, DARK, True)
    stat_card(slide, Inches(8.72), Inches(2.52), Inches(1.55), Inches(1.0), "Muddat", "Qisqaradi", SOFT_BG, GREEN)
    stat_card(slide, Inches(10.42), Inches(2.52), Inches(1.55), Inches(1.0), "Sifat", "Barqaror", SOFT_BG, BLUE)
    stat_card(slide, Inches(8.72), Inches(3.72), Inches(1.55), Inches(1.0), "Xatolik", "Kamayadi", SOFT_BG, TEAL)
    stat_card(slide, Inches(10.42), Inches(3.72), Inches(1.55), Inches(1.0), "Nazorat", "Kuchli", SOFT_BG, GREEN_DARK)
    add_textbox(
        slide, Inches(8.72), Inches(4.95), Inches(3.2), Inches(0.95),
        "Natijada loyiha qarorlari faktlarga asoslangan bo‘ladi, manfaatdor tomonlar esa yakuniy yechimni aniqroq tasavvur qiladi.",
        12, GRAY
    )
    add_footer(slide, 2)


def slide_3():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "2. Avtomatlashtirilgan loyihalashning asosiy tarkibi", "Texnologik poydevor: model, parametr, hisoblash va muvofiqlashtirish")

    process_node(slide, Inches(0.9), Inches(1.75), Inches(2.45), Inches(1.95), "BIM model", "Binoning geometriyasi, materiallari, element xossalari va hujjatlari yagona raqamli muhitda birlashtiriladi.", GREEN)
    process_node(slide, Inches(3.62), Inches(1.75), Inches(2.45), Inches(1.95), "Parametrik bog‘lanish", "O‘lcham, modul, qavat balandligi yoki element turi o‘zgarsa, bog‘liq qismlar avtomatik qayta hisoblanadi.", BLUE)
    process_node(slide, Inches(6.34), Inches(1.75), Inches(2.45), Inches(1.95), "Analitik tahlil", "Konstruktiv, energiya, hajm va muddat ko‘rsatkichlari loyihaning erta bosqichida baholanadi.", TEAL)
    process_node(slide, Inches(9.06), Inches(1.75), Inches(2.45), Inches(1.95), "Hamkorlik", "Arxitektor, konstruktor, MEP muhandis va boshqaruv guruhi bir xil ma’lumot asosida ishlaydi.", GREEN_DARK)

    for x1, x2 in [(Inches(3.35), Inches(3.62)), (Inches(6.07), Inches(6.34)), (Inches(8.79), Inches(9.06))]:
        add_line(slide, x1, Inches(2.72), x2, Inches(2.72), MID_GRAY, 1.6)

    add_rect(slide, Inches(1.0), Inches(4.25), Inches(11.5), Inches(1.85), SOFT_BG, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(1.25), Inches(4.48), Inches(2.3), Inches(0.3), "Nega aynan shu tarkib muhim?", 17, DARK, True)
    add_textbox(
        slide, Inches(1.25), Inches(4.86), Inches(10.8), Inches(0.95),
        "Avtomatlashtirilgan qurilish loyihalash faqat chizma chizishni tezlatmaydi. U loyiha ma’lumotlarini yashovchan tizimga aylantiradi: elementlar o‘zaro bog‘lanadi, o‘zgarishlar izchil tarqaladi, hisob-kitoblar dalillanadi va manfaatdor tomonlar bir xil ma’lumot manbasi bilan ishlaydi. Shu sababli qarorlar sifati, muvofiqlashtirish darajasi va yakuniy bajarilish ishonchliligi sezilarli ravishda oshadi.",
        13, GRAY
    )

    stat_card(slide, Inches(1.1), Inches(6.28), Inches(2.55), Inches(0.6), "Yagona haqiqat manbasi", "Model markazda")
    stat_card(slide, Inches(3.92), Inches(6.28), Inches(2.55), Inches(0.6), "Tez yangilanish", "Parametrlar orqali")
    stat_card(slide, Inches(6.74), Inches(6.28), Inches(2.55), Inches(0.6), "Tahliliy qaror", "Ma’lumot asosida")
    stat_card(slide, Inches(9.56), Inches(6.28), Inches(2.55), Inches(0.6), "Jamoaviy ish", "Raqamli integratsiya")
    add_footer(slide, 3)


def slide_4():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "3. Jarayon bosqichlari", "Avtomatlashtirilgan loyiha ishlab chiqishning amaliy ketma-ketligi")

    nodes = [
        ("Talablarni yig‘ish", "Buyurtmachi maqsadi, maydon cheklovlari, normativ va byudjet shartlari aniq shakllantiriladi."),
        ("Raqamli model tuzish", "Arxitektura, konstruksiya va muhandislik qatlamlari asosiy parametrlar bilan yaratiladi."),
        ("Avtomatik tekshiruv", "To‘qnashuvlar, me’yoriy cheklovlar va ma’lumot to‘liqligi bo‘yicha nazorat bajariladi."),
        ("Tahlil va optimallashtirish", "Energiya, material sarfi, konstruktiv yechim va muddat ssenariylari solishtiriladi."),
        ("Hujjatlashtirish", "Ishchi chizmalar, spesifikatsiyalar va hisobotlar modeldan avtomatik hosil qilinadi."),
        ("Monitoring va yangilash", "Qurilish davomida o‘zgarishlar modelga qayta kiritilib, loyiha bazasi dolzarb saqlanadi.")
    ]

    positions = [
        (Inches(0.95), Inches(1.7)), (Inches(4.55), Inches(1.7)), (Inches(8.15), Inches(1.7)),
        (Inches(0.95), Inches(4.15)), (Inches(4.55), Inches(4.15)), (Inches(8.15), Inches(4.15))
    ]

    for i, ((title, body), (x, y)) in enumerate(zip(nodes, positions), start=1):
        add_rect(slide, x, y, Inches(3.0), Inches(1.75), WHITE, LIGHT_GRAY, 1.0, radius=True)
        draw_circle(slide, x + Inches(0.38), y + Inches(0.38), Inches(0.42), GREEN if i % 2 else BLUE, WHITE, 1.0)
        add_textbox(slide, x + Inches(0.27), y + Inches(0.24), Inches(0.22), Inches(0.18), str(i), 12, WHITE, True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.68), y + Inches(0.18), Inches(2.0), Inches(0.28), title, 15, DARK, True)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.58), Inches(2.55), Inches(0.95), body, 11, GRAY)

    add_line(slide, Inches(3.95), Inches(2.58), Inches(4.55), Inches(2.58), MID_GRAY, 1.4)
    add_line(slide, Inches(7.55), Inches(2.58), Inches(8.15), Inches(2.58), MID_GRAY, 1.4)
    add_line(slide, Inches(9.65), Inches(3.45), Inches(9.65), Inches(4.15), MID_GRAY, 1.4)
    add_line(slide, Inches(8.15), Inches(5.03), Inches(7.55), Inches(5.03), MID_GRAY, 1.4)
    add_line(slide, Inches(4.55), Inches(5.03), Inches(3.95), Inches(5.03), MID_GRAY, 1.4)

    add_footer(slide, 4)


def slide_5():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "4. Parametrik yondashuv va BIM afzalligi", "Model ichidagi bog‘lanishlar loyihadagi o‘zgarishlarni boshqarishni yengillashtiradi")

    add_rect(slide, Inches(0.82), Inches(1.62), Inches(5.25), Inches(4.95), SOFT_BG, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(1.05), Inches(1.92), Inches(4.4), Inches(0.35), "Parametrik yondashuv qanday ishlaydi?", 17, DARK, True)
    add_textbox(
        slide, Inches(1.05), Inches(2.32), Inches(4.62), Inches(2.2),
        "Parametrik loyihalashda devor qalinligi, modul o‘lchami, qavat balandligi, ustun oralig‘i yoki fasad panellari soni kabi o‘zgaruvchilar oldindan mantiqiy bog‘lanadi. Natijada loyiha ichidagi bir parametr o‘zgarganda, u bilan bog‘liq ko‘plab elementlar avtomatik yangilanadi. Bu ayniqsa takrorlanuvchi qavatlar, yirik majmualar va bir necha variantli konsepsiyalarni ishlab chiqishda juda samarali.",
        13, GRAY
    )
    add_textbox(
        slide, Inches(1.05), Inches(4.72), Inches(4.62), Inches(1.45),
        "BIM esa ushbu parametrlarni faqat geometriya bilan cheklab qo‘ymaydi; elementning materiali, markasi, hajmi, texnik tavsifi va ekspluatatsiya ma’lumotlari ham model ichiga biriktiriladi. Shuning uchun model chizma manbai bo‘lish bilan birga boshqaruv manbaiga ham aylanadi.",
        13, GRAY
    )

    add_rect(slide, Inches(6.35), Inches(1.62), Inches(6.1), Inches(2.2), WHITE, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(6.62), Inches(1.9), Inches(2.7), Inches(0.3), "Amaliy ustunliklar", 17, DARK, True)
    add_paragraphs_box(
        slide, Inches(6.55), Inches(2.28), Inches(5.4), Inches(1.18),
        [
            "Bir necha loyiha variantini tez ishlab chiqish va solishtirish mumkin.",
            "Reja, fasad, kesim va jadval o‘rtasidagi moslik yuqori darajada saqlanadi.",
            "O‘zgarish kiritish narxi va vaqt sarfi kamayadi.",
            "Katta hajmdagi loyihalarda sifat nazorati yengillashadi."
        ],
        12, GRAY
    )

    add_rect(slide, Inches(6.35), Inches(4.05), Inches(6.1), Inches(2.52), WHITE, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(6.62), Inches(4.32), Inches(2.4), Inches(0.3), "Oddiy misol", 17, DARK, True)

    stat_card(slide, Inches(6.65), Inches(4.82), Inches(1.7), Inches(1.1), "Qavat balandligi", "+0.3 m", SOFT_BG, GREEN)
    stat_card(slide, Inches(8.48), Inches(4.82), Inches(1.7), Inches(1.1), "Kesimlar", "Yangilanadi", SOFT_BG, BLUE)
    stat_card(slide, Inches(10.31), Inches(4.82), Inches(1.7), Inches(1.1), "Spetsifikatsiya", "Qayta hisob", SOFT_BG, TEAL)

    add_textbox(
        slide, Inches(6.65), Inches(6.02), Inches(5.35), Inches(0.4),
        "Ya’ni bitta o‘zgarish hujjatlar zanjiri bo‘ylab nazoratli ravishda tarqaladi va qo‘lda tahrir hajmini kamaytiradi.",
        11, GRAY
    )
    add_footer(slide, 5)


def slide_6():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "5. Doirasimon diagramma: funksional taqsimot", "Avtomatlashtirilgan loyihalashdagi asosiy yo‘nalishlar ulushi")

    data = [
        ("BIM modellashtirish", 35),
        ("Parametrik chizmalar", 25),
        ("Analitik hisoblash", 20),
        ("Nazorat va integratsiya", 20),
    ]
    colors = [GREEN, BLUE, TEAL, GREEN_DARK]

    add_rect(slide, Inches(0.82), Inches(1.55), Inches(12.0), Inches(5.45), WHITE, LIGHT_GRAY, 1.0, radius=True)
    donut_chart_slide(slide, Inches(1.15), Inches(2.0), Inches(3.2), data, colors, "100%", "Yagona ekotizim")

    add_rect(slide, Inches(8.4), Inches(1.98), Inches(3.85), Inches(4.4), SOFT_BG, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(8.68), Inches(2.25), Inches(2.8), Inches(0.3), "Diagrammadan xulosa", 17, DARK, True)
    add_textbox(
        slide, Inches(8.68), Inches(2.65), Inches(3.05), Inches(2.95),
        "Doirasimon taqsimotdan ko‘rinadiki, eng katta ulush BIM modellashtirishga to‘g‘ri keladi, chunki aynan model barcha ma’lumot oqimining markazidir. Parametrik chizmalar esa loyiha o‘zgarishlarini tez boshqarish uchun ikkinchi muhim qatlam hisoblanadi. Analitik hisoblash va nazorat-integratsiya yo‘nalishlari ulushi nisbatan teng bo‘lsa-da, ular qarorlar sifatini oshirish hamda jamoalararo ishlashni muvofiqlashtirishda hal qiluvchi rol o‘ynaydi.",
        12, GRAY
    )
    add_footer(slide, 6)


def slide_7():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "6. Loyiha boshqaruvida qo‘llanilishi", "Avtomatlashtirish nafaqat chizma, balki boshqaruv vositasi hamdir")

    add_rect(slide, Inches(0.85), Inches(1.68), Inches(3.85), Inches(4.95), SOFT_BG, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(1.12), Inches(1.98), Inches(2.9), Inches(0.3), "Boshqaruv nuqtalari", 17, DARK, True)
    add_paragraphs_box(
        slide, Inches(1.08), Inches(2.35), Inches(3.25), Inches(3.9),
        [
            "Vazifalarni loyiha bosqichlari kesimida taqsimlash va mas’ullarni biriktirish osonlashadi.",
            "Versiyalarni boshqarish orqali qaysi model yoki hujjat dolzarb ekanligi aniq ko‘rinadi.",
            "Tuzatishlar tarixini kuzatish va tasdiqlash zanjirini yuritish soddalashadi.",
            "Muddat, hajm va material bo‘yicha hisobotlar tez tayyorlanadi.",
            "Qurilish maydoniga uzatiladigan hujjatlar izchil va bir manbadan shakllanadi."
        ],
        12, GRAY
    )

    add_rect(slide, Inches(4.98), Inches(1.68), Inches(3.0), Inches(4.95), WHITE, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(5.25), Inches(1.98), Inches(2.2), Inches(0.3), "Kimlar foyda ko‘radi?", 17, DARK, True)
    process_node(slide, Inches(5.2), Inches(2.42), Inches(2.55), Inches(0.92), "Buyurtmachi", "Loyiha holatini tushunarli ko‘radi.", GREEN)
    process_node(slide, Inches(5.2), Inches(3.45), Inches(2.55), Inches(0.92), "Loyiha jamoasi", "O‘zgarishlarni tez moslashtiradi.", BLUE)
    process_node(slide, Inches(5.2), Inches(4.48), Inches(2.55), Inches(0.92), "Pudratchi", "Aniqroq ishchi hujjat oladi.", TEAL)
    process_node(slide, Inches(5.2), Inches(5.51), Inches(2.55), Inches(0.92), "Ekspluatator", "Keyingi foydalanish ma’lumotiga ega bo‘ladi.", GREEN_DARK)

    add_rect(slide, Inches(8.25), Inches(1.68), Inches(4.25), Inches(4.95), SOFT_BG, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(8.55), Inches(1.98), Inches(3.2), Inches(0.3), "Boshqaruvdagi strategik ta’sir", 17, DARK, True)
    add_textbox(
        slide, Inches(8.55), Inches(2.38), Inches(3.4), Inches(3.8),
        "Avtomatlashtirilgan loyihalash boshqaruv guruhiga faqat tayyor chizmalarni emas, balki qaror chiqarish uchun kerakli ma’lumotlarni ham beradi. Qaysi element qimmatlashgani, qayerda to‘qnashuv yuzaga kelgani, qaysi zona kechikishga moyilligi yoki qaysi variant resursni kamroq talab qilishi tez aniqlanadi. Shu tufayli loyiha rahbariyati reaktiv emas, balki proaktiv usulda boshqaruv olib boradi.",
        12, GRAY
    )
    add_footer(slide, 7)


def slide_8():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "7. Xavflar va joriy etishdagi muammolar", "Har qanday transformatsiya kabi bu yo‘nalishda ham tayyorgarlik muhim")

    add_rect(slide, Inches(0.88), Inches(1.72), Inches(12.0), Inches(4.95), WHITE, LIGHT_GRAY, 1.0, radius=True)

    cols = [
        ("Tashkiliy muammolar", [
            "Jamoa ichida raqamli madaniyat yetarli bo‘lmasligi mumkin.",
            "Ish jarayonlari eski odatlarga bog‘langan bo‘lsa, yangi tartibni qabul qilish sekinlashadi.",
            "Mas’uliyat va ma’lumot egasini belgilash aniq bo‘lmasa, model sifati pasayadi."
        ], GREEN),
        ("Texnik muammolar", [
            "Shablon, standart va nomlash qoidalari bo‘lmasa, loyiha fayllari tartibsizlashadi.",
            "Og‘ir modellar kompyuter resurslariga yuqori talab qo‘yadi.",
            "Turli bo‘limlarning ma’lumot almashinuvi yetarli darajada sozlanmasa, integratsiya qiyinlashadi."
        ], BLUE),
        ("Menejment xatarlari", [
            "Boshlang‘ich bosqichda ta’lim va tizimlashtirish uchun qo‘shimcha vaqt talab etiladi.",
            "Investitsiya samarasi darhol emas, bosqichma-bosqich ko‘rinadi.",
            "Loyihani avtomatlashtirishni faqat dastur xaridi deb talqin qilish noto‘g‘ri natija beradi."
        ], TEAL),
    ]

    x_positions = [Inches(1.05), Inches(4.35), Inches(7.65)]
    for (title, items, color), x in zip(cols, x_positions):
        add_rect(slide, x, Inches(2.0), Inches(2.95), Inches(4.3), SOFT_BG, LIGHT_GRAY, 1.0, radius=True)
        add_rect(slide, x, Inches(2.0), Inches(2.95), Inches(0.18), color, color, 0.5, radius=True)
        add_textbox(slide, x + Inches(0.22), Inches(2.22), Inches(2.45), Inches(0.3), title, 16, DARK, True)
        add_paragraphs_box(slide, x + Inches(0.18), Inches(2.62), Inches(2.5), Inches(3.2), items, 11, GRAY)

    add_rect(slide, Inches(10.95), Inches(2.0), Inches(1.3), Inches(4.3), WHITE, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(11.13), Inches(2.24), Inches(0.9), Inches(0.3), "Yechim", 15, DARK, True, align=PP_ALIGN.CENTER)
    draw_circle(slide, Inches(11.6), Inches(3.2), Inches(0.72), GREEN_LIGHT, WHITE, 1.0)
    add_textbox(slide, Inches(11.22), Inches(3.0), Inches(0.76), Inches(0.36), "1", 20, GREEN_DARK, True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(11.08), Inches(3.55), Inches(1.0), Inches(0.55), "Standart", 11, GRAY, align=PP_ALIGN.CENTER)
    draw_circle(slide, Inches(11.6), Inches(4.35), Inches(0.72), BLUE_LIGHT, WHITE, 1.0)
    add_textbox(slide, Inches(11.22), Inches(4.15), Inches(0.76), Inches(0.36), "2", 20, BLUE, True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(11.05), Inches(4.7), Inches(1.1), Inches(0.55), "Ta’lim", 11, GRAY, align=PP_ALIGN.CENTER)
    draw_circle(slide, Inches(11.6), Inches(5.5), Inches(0.72), GREEN_LIGHT, WHITE, 1.0)
    add_textbox(slide, Inches(11.22), Inches(5.3), Inches(0.76), Inches(0.36), "3", 20, GREEN_DARK, True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(11.0), Inches(5.85), Inches(1.2), Inches(0.55), "Bosqichli joriy etish", 10, GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)


def slide_9():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title(slide, "8. Joriy etish bo‘yicha tavsiyalar", "Amaliy yo‘l xaritasi: kichikdan boshlash, standartlashtirish va kengaytirish")

    add_rect(slide, Inches(0.92), Inches(1.72), Inches(11.85), Inches(5.0), WHITE, LIGHT_GRAY, 1.0, radius=True)

    steps = [
        ("1-bosqich", "Pilot loyiha tanlash", "Murakkabligi boshqariladigan, lekin real qiymat ko‘rsata oladigan loyiha tanlanadi."),
        ("2-bosqich", "Shablon va standart", "Element nomlari, qatlamlar, oila/family tamoyillari, jadval va hujjat formatlari belgilanadi."),
        ("3-bosqich", "Jamoani tayyorlash", "Arxitektor, konstruktor va muhandislar uchun rollarga mos treninglar tashkil etiladi."),
        ("4-bosqich", "Nazorat ssenariysi", "To‘qnashuv, to‘liqlik va sifat bo‘yicha davriy tekshiruv qoidalari ishlab chiqiladi."),
        ("5-bosqich", "Kengaytirish", "Pilotdan olingan tajriba asosida yechim boshqa loyihalarga bosqichma-bosqich yoyiladi.")
    ]

    y = Inches(2.0)
    for i, (phase, title, body) in enumerate(steps):
        add_rect(slide, Inches(1.15), y, Inches(10.9), Inches(0.8), SOFT_BG if i % 2 == 0 else WHITE, LIGHT_GRAY, 1.0, radius=True)
        add_rect(slide, Inches(1.3), y + Inches(0.14), Inches(1.1), Inches(0.5), GREEN if i % 2 == 0 else BLUE, GREEN if i % 2 == 0 else BLUE, 0.5, radius=True)
        add_textbox(slide, Inches(1.45), y + Inches(0.22), Inches(0.8), Inches(0.2), phase, 11, WHITE, True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(2.65), y + Inches(0.12), Inches(2.15), Inches(0.25), title, 15, DARK, True)
        add_textbox(slide, Inches(4.95), y + Inches(0.12), Inches(6.65), Inches(0.42), body, 11, GRAY)
        y += Inches(0.92)

    add_textbox(
        slide, Inches(1.18), Inches(6.72), Inches(10.8), Inches(0.3),
        "Eng muhim qoida: avtomatlashtirishni faqat dasturiy vosita sifatida emas, balki standart, jarayon va malaka uyg‘unligi sifatida joriy etish kerak.",
        12, GRAY
    )
    add_footer(slide, 9)


def slide_10():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)

    add_rect(slide, Inches(0.68), Inches(0.68), Inches(12.0), Inches(5.95), WHITE, LIGHT_GRAY, 1.2, radius=True)
    add_rect(slide, Inches(0.9), Inches(0.9), Inches(11.56), Inches(0.22), GREEN, GREEN, 0.5, radius=True)
    add_textbox(slide, Inches(1.05), Inches(1.35), Inches(7.0), Inches(0.7), "Xulosa va yakuniy fikrlar", 27, DARK, True)

    add_textbox(
        slide, Inches(1.08), Inches(2.05), Inches(6.7), Inches(2.7),
        "Avtomatlashtirilgan qurilish loyihalash bugungi kunda samaradorlik, sifat va boshqaruv shaffofligini oshiradigan muhim yondashuvga aylandi. Uning kuchi BIM model, parametrik bog‘lanish, avtomatik nazorat va analitik tahlilning yagona jarayonda ishlashidadir. Natijada loyiha hujjatlari tezroq tayyorlanadi, xatolar erta aniqlanadi, jamoalar o‘rtasidagi kelishuv kuchayadi va buyurtmachiga ko‘proq asoslangan qarorlar taklif qilinadi.",
        14, GRAY
    )
    add_textbox(
        slide, Inches(1.08), Inches(4.85), Inches(6.7), Inches(1.25),
        "Kelajakda ushbu yondashuv faqat loyihalash bosqichidagina emas, balki qurilish, ekspluatatsiya va aktivlarni boshqarish tizimlari bilan yanada chuqur integratsiyalashadi. Demak, avtomatlashtirish — bu tanlov emas, raqobatbardoshlik omili.",
        14, GRAY
    )

    add_rect(slide, Inches(8.25), Inches(1.52), Inches(3.75), Inches(4.7), SOFT_BG, LIGHT_GRAY, 1.0, radius=True)
    add_textbox(slide, Inches(8.55), Inches(1.82), Inches(2.7), Inches(0.3), "Yakuniy tezislar", 18, DARK, True)
    add_paragraphs_box(
        slide, Inches(8.48), Inches(2.25), Inches(3.0), Inches(2.7),
        [
            "Yagona model qarorlar sifatini oshiradi.",
            "Parametrik yondashuv o‘zgarishlarni tez boshqaradi.",
            "Analitik tekshiruv resurslarni tejaydi.",
            "Standart va trening joriy etish muvaffaqiyatini belgilaydi.",
            "Bosqichma-bosqich tatbiq etish eng amaliy yo‘ldir."
        ],
        12, GRAY
    )
    stat_card(slide, Inches(8.55), Inches(5.15), Inches(1.45), Inches(0.82), "Natija", "Tezroq", SOFT_BG, GREEN)
    stat_card(slide, Inches(10.1), Inches(5.15), Inches(1.45), Inches(0.82), "Sifat", "Yuqori", SOFT_BG, BLUE)

    add_footer(slide, 10)


slide_1()
slide_2()
slide_3()
slide_4()
slide_5()
slide_6()
slide_7()
slide_8()
slide_9()
slide_10()

while len(prs.slides) > SLIDE_COUNT:
    rId = prs.slides._sldIdLst[-1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[-1]

prs.save(OUTPUT_PATH)