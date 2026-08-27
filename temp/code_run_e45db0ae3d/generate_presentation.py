import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor

MAVZU = "Avtomatlashtirilgan zamonaviy bino qurilish tehnikalari"
SLIDE_COUNT = 10
TAQDIMOT_TILI = "o‘zbek"
TAQDIMOT_EGASI = "Javlonbek Moʻydinov"
OUTPUT_PATH = os.environ.get("PPTX_OUTPUT_PATH", "presentation.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(39, 52, 64)
GRAY = RGBColor(108, 117, 125)
LIGHT_GRAY = RGBColor(230, 234, 238)
MID_GRAY = RGBColor(200, 206, 212)
GREEN = RGBColor(34, 139, 94)
GREEN_DARK = RGBColor(22, 110, 73)
GREEN_LIGHT = RGBColor(224, 245, 236)
BLUE = RGBColor(40, 98, 255)
BLUE_DARK = RGBColor(22, 64, 143)
BLUE_LIGHT = RGBColor(232, 240, 255)
SOFT = RGBColor(245, 247, 249)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"

def set_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=1.2, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.color.rgb = line_color if line_color else fill_color
    shp.line.width = Pt(line_width)
    return shp

def add_line(slide, x1, y1, x2, y2, color=MID_GRAY, width=1.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line

def add_textbox(slide, x, y, w, h, text, font_size=20, color=DARK, bold=False, font_name=BODY_FONT,
                align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP, margin=0.05, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = color
    return tb

def add_paragraphs_box(slide, x, y, w, h, paragraphs, font_size=18, color=DARK, bullet=True,
                       bg=None, line=None, radius=True, title=None):
    if bg is not None:
        add_rect(slide, x, y, w, h, bg, line if line else bg, 1.0, radius=radius)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    if title:
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.LEFT
        r0 = p0.add_run()
        r0.text = title
        r0.font.name = TITLE_FONT
        r0.font.size = Pt(font_size + 2)
        r0.font.bold = True
        r0.font.color.rgb = DARK
        for item in paragraphs:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.level = 0
            p.bullet = bullet
            r = p.add_run()
            r.text = item
            r.font.name = BODY_FONT
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
    else:
        first = True
        for item in paragraphs:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.level = 0
            p.bullet = bullet
            r = p.add_run()
            r.text = item
            r.font.name = BODY_FONT
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
    return tb

def add_title(slide, title, subtitle=None, accent="left"):
    if accent == "left":
        add_rect(slide, Inches(0.45), Inches(0.45), Inches(0.16), Inches(1.05), GREEN, GREEN, 0.5, radius=True)
        add_textbox(slide, Inches(0.8), Inches(0.42), Inches(8.8), Inches(0.72), title, 28, DARK, True, TITLE_FONT)
        if subtitle:
            add_textbox(slide, Inches(0.82), Inches(1.07), Inches(8.5), Inches(0.42), subtitle, 11.5, GRAY, False, BODY_FONT)
        add_line(slide, Inches(0.8), Inches(1.48), Inches(12.7), Inches(1.48), LIGHT_GRAY, 1.2)
    elif accent == "band":
        add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), SOFT, SOFT, 0.5, radius=False)
        add_rect(slide, Inches(0.45), Inches(0.32), Inches(2.6), Inches(0.12), GREEN, GREEN, 0.1, radius=True)
        add_textbox(slide, Inches(0.52), Inches(0.18), Inches(9.8), Inches(0.42), title, 26, DARK, True, TITLE_FONT)
        if subtitle:
            add_textbox(slide, Inches(0.55), Inches(0.58), Inches(8.8), Inches(0.26), subtitle, 11.5, GRAY, False, BODY_FONT)
    elif accent == "topline":
        add_line(slide, Inches(0.4), Inches(0.42), Inches(12.95), Inches(0.42), GREEN, 2.4)
        add_textbox(slide, Inches(0.45), Inches(0.56), Inches(9.2), Inches(0.55), title, 27, DARK, True, TITLE_FONT)
        if subtitle:
            add_textbox(slide, Inches(0.48), Inches(1.05), Inches(8.8), Inches(0.32), subtitle, 11.5, GRAY, False, BODY_FONT)

def add_footer(slide, idx):
    add_line(slide, Inches(0.5), Inches(7.05), Inches(12.8), Inches(7.05), LIGHT_GRAY, 1.0)
    add_textbox(slide, Inches(0.55), Inches(7.08), Inches(4.2), Inches(0.22), TAQDIMOT_EGASI, 9.5, GRAY, False, BODY_FONT)
    add_textbox(slide, Inches(11.9), Inches(7.08), Inches(0.7), Inches(0.22), str(idx), 9.5, GRAY, False, BODY_FONT, PP_ALIGN.RIGHT)

def add_badge(slide, x, y, w, h, text, fill=GREEN_LIGHT, line=GREEN, color=GREEN_DARK):
    shp = add_rect(slide, x, y, w, h, fill, line, 1.0, radius=True)
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = BODY_FONT
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = color
    return shp

def add_icon_house(slide, x, y, scale=1.0, color=GREEN):
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, x, y, Inches(0.5)*scale, Inches(0.34)*scale).fill.solid()
    roof = slide.shapes._spTree[-1]
    roof.fill.fore_color.rgb = color
    roof.line.color.rgb = color
    body = add_rect(slide, x + Inches(0.08)*scale, y + Inches(0.28)*scale, Inches(0.34)*scale, Inches(0.28)*scale, GREEN_LIGHT, color, 1.0, radius=False)
    door = add_rect(slide, x + Inches(0.21)*scale, y + Inches(0.38)*scale, Inches(0.07)*scale, Inches(0.18)*scale, color, color, 1.0, radius=False)
    return body, door

def add_icon_robot(slide, x, y, scale=1.0):
    head = add_rect(slide, x, y, Inches(0.34)*scale, Inches(0.24)*scale, BLUE_LIGHT, BLUE, 1.0, radius=True)
    eye1 = add_rect(slide, x + Inches(0.08)*scale, y + Inches(0.08)*scale, Inches(0.04)*scale, Inches(0.04)*scale, BLUE, BLUE, 0.5, radius=True)
    eye2 = add_rect(slide, x + Inches(0.22)*scale, y + Inches(0.08)*scale, Inches(0.04)*scale, Inches(0.04)*scale, BLUE, BLUE, 0.5, radius=True)
    body = add_rect(slide, x + Inches(0.05)*scale, y + Inches(0.26)*scale, Inches(0.24)*scale, Inches(0.22)*scale, GREEN_LIGHT, GREEN, 1.0, radius=True)
    add_line(slide, x + Inches(0.17)*scale, y - Inches(0.04)*scale, x + Inches(0.17)*scale, y, BLUE, 1.0)
    return head, body, eye1, eye2

def add_stat_card(slide, x, y, w, h, number, label, fill, line, num_color=DARK):
    shp = add_rect(slide, x, y, w, h, fill, line, 1.0, radius=True)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.12), w - Inches(0.24), Inches(0.42), number, 24, num_color, True, TITLE_FONT)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.58), w - Inches(0.24), h - Inches(0.68), label, 11.5, GRAY, False, BODY_FONT)
    return shp

def add_process_box(slide, x, y, w, h, title, text, num, fill):
    add_rect(slide, x, y, w, h, fill, MID_GRAY, 0.8, radius=True)
    add_badge(slide, x + Inches(0.12), y + Inches(0.12), Inches(0.55), Inches(0.32), num, GREEN_LIGHT, GREEN, GREEN_DARK)
    add_textbox(slide, x + Inches(0.78), y + Inches(0.1), w - Inches(0.9), Inches(0.36), title, 15, DARK, True, TITLE_FONT)
    add_textbox(slide, x + Inches(0.12), y + Inches(0.48), w - Inches(0.24), h - Inches(0.6), text, 11.5, DARK, False, BODY_FONT)
    return True

def add_donut(slide, cx, cy, size, values, colors, labels):
    total = float(sum(values))
    start_angle = -90.0
    outer_r = size / 2.0
    inner_r = outer_r * 0.58
    center_x = cx + outer_r
    center_y = cy + outer_r

    for val, color in zip(values, colors):
        angle = 360.0 * val / total
        seg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.PIE, cx, cy, size, size)
        seg.fill.solid()
        seg.fill.fore_color.rgb = color
        seg.line.color.rgb = WHITE
        seg.line.width = Pt(1.0)
        seg.adjustments[0] = int(start_angle * 60000)
        seg.adjustments[1] = int((start_angle + angle) * 60000)
        start_angle += angle

    inner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, center_x - inner_r, center_y - inner_r, inner_r*2, inner_r*2)
    inner.fill.solid()
    inner.fill.fore_color.rgb = WHITE
    inner.line.color.rgb = WHITE

    add_textbox(slide, cx + size*0.27, cy + size*0.36, size*0.46, size*0.18, "Ulashuv", 17, DARK, True, TITLE_FONT, PP_ALIGN.CENTER)
    add_textbox(slide, cx + size*0.22, cy + size*0.53, size*0.56, size*0.12, "texnologiya yo‘nalishlari", 9.5, GRAY, False, BODY_FONT, PP_ALIGN.CENTER)

    lx = cx + size + Inches(0.25)
    ly = cy + Inches(0.18)
    for i, (label, val, color) in enumerate(zip(labels, values, colors)):
        add_rect(slide, lx, ly + Inches(0.48)*i, Inches(0.18), Inches(0.18), color, color, 0.5, radius=False)
        pct = f"{int(round(val / total * 100))}%"
        add_textbox(slide, lx + Inches(0.25), ly + Inches(0.02) + Inches(0.48)*i, Inches(2.2), Inches(0.22),
                    f"{label} — {pct}", 11.5, DARK, False, BODY_FONT)

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_rect(slide, Inches(0.55), Inches(0.55), Inches(12.2), Inches(6.15), WHITE, LIGHT_GRAY, 1.2, radius=True)
add_rect(slide, Inches(0.75), Inches(0.85), Inches(0.18), Inches(1.25), GREEN, GREEN, 0.2, radius=True)
add_rect(slide, Inches(0.98), Inches(0.85), Inches(0.09), Inches(0.78), BLUE, BLUE, 0.2, radius=True)
add_textbox(slide, Inches(1.18), Inches(0.92), Inches(7.1), Inches(1.2), MAVZU, 28, DARK, True, TITLE_FONT)
add_textbox(slide, Inches(1.2), Inches(2.05), Inches(6.2), Inches(0.85),
            "Zamonaviy qurilish maydonlarida robototexnika, BIM, IoT sensorlar, raqamli nazorat va aqlli boshqaruv yechimlarining amaliy qo‘llanilishi.",
            16, GRAY, False, BODY_FONT)
add_badge(slide, Inches(1.2), Inches(3.0), Inches(2.0), Inches(0.38), "Premium tahliliy taqdimot", BLUE_LIGHT, BLUE, BLUE_DARK)

right = add_rect(slide, Inches(8.35), Inches(1.0), Inches(3.8), Inches(4.9), SOFT, LIGHT_GRAY, 1.0, radius=True)
add_rect(slide, Inches(8.65), Inches(1.42), Inches(3.2), Inches(2.05), WHITE, MID_GRAY, 1.0, radius=True)
add_rect(slide, Inches(9.0), Inches(2.0), Inches(2.5), Inches(1.2), GREEN_LIGHT, GREEN, 1.0, radius=True)
add_icon_house(slide, Inches(9.15), Inches(1.74), 2.8, GREEN)
add_icon_robot(slide, Inches(10.75), Inches(2.02), 2.3)
add_line(slide, Inches(9.1), Inches(3.95), Inches(11.45), Inches(3.95), MID_GRAY, 1.0)
add_textbox(slide, Inches(8.75), Inches(4.15), Inches(2.9), Inches(0.75),
            "Aqlli qurilish boshqaruvi\n• tezlik • aniqlik • xavfsizlik", 13, DARK, True, BODY_FONT)

add_textbox(slide, Inches(1.2), Inches(5.35), Inches(3.6), Inches(0.35), f"Taqdimot egasi: {TAQDIMOT_EGASI}", 13, DARK, True, BODY_FONT)
add_textbox(slide, Inches(1.2), Inches(5.7), Inches(3.8), Inches(0.28), f"Tili: {TAQDIMOT_TILI}", 11.5, GRAY, False, BODY_FONT)
add_textbox(slide, Inches(1.2), Inches(6.02), Inches(4.8), Inches(0.28), "Slaydlar soni: 10", 11.5, GRAY, False, BODY_FONT)
add_footer(slide, 1)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Avtomatlashtirilgan qurilish tushunchasi", "Asosiy ta’riflar, tizimlar va bozor drayverlari", "band")
add_paragraphs_box(
    slide, Inches(0.6), Inches(1.5), Inches(6.0), Inches(4.9),
    [
        "Avtomatlashtirilgan qurilish — bu loyiha, material, texnika, ishchi kuchi va nazorat jarayonlarini raqamli platformalar orqali uyg‘un boshqarish yondashuvidir.",
        "Mazkur modelda BIM, ERP, GPS kuzatuv, sensorli monitoring, dron kuzatuvi va robotlashtirilgan uskunalar bir-biri bilan ma’lumot almashadi.",
        "An’anaviy usuldan farqli ravishda qarorlar tajriba asosida emas, balki real vaqt ma’lumotlari, algoritmik tahlil va standartlashtirilgan ish jarayonlari bilan qabul qilinadi.",
        "Natijada qurilish sikli qisqaradi, xatoliklar kamayadi, resurs sarfi nazorat qilinadi va loyiha topshirish sifati barqarorlashadi."
    ],
    font_size=14, color=DARK, bullet=True, bg=SOFT, line=LIGHT_GRAY, radius=True, title="Mazmuni"
)
add_rect(slide, Inches(6.9), Inches(1.7), Inches(5.75), Inches(4.55), WHITE, LIGHT_GRAY, 1.0, radius=True)
add_stat_card(slide, Inches(7.2), Inches(2.05), Inches(1.7), Inches(1.15), "24/7", "Avtomatik monitoring va qurilish maydoni kuzatuvi", GREEN_LIGHT, GREEN)
add_stat_card(slide, Inches(9.0), Inches(2.05), Inches(1.7), Inches(1.15), "±", "Aniqlikni oshiruvchi sensor va lazer o‘lchovi", BLUE_LIGHT, BLUE)
add_stat_card(slide, Inches(10.8), Inches(2.05), Inches(1.55), Inches(1.15), "BIM", "Yagona raqamli model va muvofiqlashtirish", SOFT, MID_GRAY)
add_textbox(slide, Inches(7.2), Inches(3.55), Inches(5.0), Inches(1.6),
            "Bozorni harakatga keltirayotgan omillar:\n"
            "• mehnat unumdorligini oshirish zarurati\n"
            "• material isrofini kamaytirish bosimi\n"
            "• murakkab ob’ektlar sonining ortishi\n"
            "• xavfsizlik va hujjatlashtirish talablarining kuchayishi",
            12.5, DARK, False, BODY_FONT)
add_badge(slide, Inches(7.2), Inches(5.35), Inches(2.35), Inches(0.34), "Raqamli qurilish ekotizimi", GREEN_LIGHT, GREEN, GREEN_DARK)
add_footer(slide, 2)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Asosiy texnologik bloklar", "Qurilishni avtomatlashtirishni tashkil etuvchi tayanch yechimlar", "topline")

cards = [
    ("BIM va raqamli model", "Loyiha, muhandislik tarmoqlari, kolliziya nazorati va smeta ma’lumotlari yagona raqamli makonda boshqariladi.", GREEN_LIGHT, GREEN),
    ("Robotlashtirilgan texnika", "Avtonom yoki yarim avtonom texnika qazish, beton quyish, kesish va tashish kabi amallarni takroriy aniqlik bilan bajaradi.", BLUE_LIGHT, BLUE),
    ("IoT sensorlar", "Harorat, namlik, vibratsiya, chang, uskunaning holati va energiya iste’moli real vaqt rejimida kuzatiladi.", SOFT, MID_GRAY),
    ("Dron va skanerlash", "Topografik nazorat, hajm hisoblash, ish hajmini foto-hisobotlash hamda xavfli hududlarni masofadan ko‘rish imkonini beradi.", GREEN_LIGHT, GREEN),
    ("Ma’lumotlar tahlili", "Ish jadvali, xarajat, kechikish va resurs bandligini tahlil qilish orqali boshqaruv qarorlari tezlashtiriladi.", BLUE_LIGHT, BLUE),
    ("Aqlli boshqaruv paneli", "Rahbariyat uchun KPI, ogohlantirishlar, progress va sifat ko‘rsatkichlari tushunarli ko‘rinishda taqdim etiladi.", SOFT, MID_GRAY),
]

x_positions = [Inches(0.65), Inches(4.45), Inches(8.25)]
y_positions = [Inches(1.55), Inches(4.35)]
idx = 0
for y in y_positions:
    for x in x_positions:
        title, text, fillc, linec = cards[idx]
        add_rect(slide, x, y, Inches(3.35), Inches(2.3), fillc, linec, 1.0, radius=True)
        add_rect(slide, x + Inches(0.15), y + Inches(0.16), Inches(0.22), Inches(0.22), linec, linec, 0.5, radius=True)
        add_textbox(slide, x + Inches(0.48), y + Inches(0.12), Inches(2.68), Inches(0.35), title, 14.5, DARK, True, TITLE_FONT)
        add_textbox(slide, x + Inches(0.16), y + Inches(0.56), Inches(2.95), Inches(1.45), text, 11.5, DARK, False, BODY_FONT)
        idx += 1
add_footer(slide, 3)

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Avtomatlashtirilgan ish jarayoni", "Loyihaning boshidan topshirishgacha bo‘lgan raqamli oqim", "left")
steps = [
    ("01", "Loyiha modellashtirish", "Arxitektura, konstruksiya va MEP bo‘limlari BIM muhitida integratsiyalanadi, kolliziyalar oldindan aniqlanadi.", GREEN_LIGHT),
    ("02", "Rejalashtirish va resurslash", "Jadval, texnika, brigada va material oqimi ishlab chiqilib, ustuvor ish frontlari belgilab olinadi.", SOFT),
    ("03", "Maydonni raqamli tayyorlash", "GPS, dron va lazer o‘lchovlari orqali maydon relefi, xavf zonalari va logistik yo‘laklar aniqlashtiriladi.", BLUE_LIGHT),
    ("04", "Avtomatlashtirilgan bajarish", "Robotlashtirilgan yoki sensorli texnika asosiy ishlarni takroriy aniqlikda bajaradi, ma’lumot tizimga uzatiladi.", GREEN_LIGHT),
    ("05", "Monitoring va sifat nazorati", "Har bir bosqich foto, sensor va chek-listlar bilan tasdiqlanadi; chetlanishlar bo‘yicha ogohlantirish beriladi.", SOFT),
]
start_x = Inches(0.7)
w = Inches(2.35)
gap = Inches(0.18)
for i, item in enumerate(steps):
    x = start_x + i * (w + gap)
    add_process_box(slide, x, Inches(2.0), w, Inches(3.8), item[1], item[2], item[0], item[3])
    if i < len(steps) - 1:
        add_line(slide, x + w, Inches(3.9), x + w + gap - Inches(0.03), Inches(3.9), GREEN, 2.0)
add_badge(slide, Inches(0.75), Inches(6.2), Inches(3.0), Inches(0.34), "Jarayonlar bir-biriga bog‘langan", BLUE_LIGHT, BLUE, BLUE_DARK)
add_footer(slide, 4)

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Doirasimon diagramma: texnologiyalar ulushi", "Qurilishda qo‘llanayotgan avtomatlashtirish yo‘nalishlari bo‘yicha shartli taqsimot", "band")
add_rect(slide, Inches(0.65), Inches(1.55), Inches(7.0), Inches(5.1), WHITE, LIGHT_GRAY, 1.0, radius=True)
add_donut(
    slide,
    Inches(1.1), Inches(2.0), Inches(3.2),
    [30, 22, 18, 16, 14],
    [GREEN, BLUE, GREEN_DARK, GRAY, BLUE_DARK],
    ["BIM", "Sensorlar", "Robot texnika", "Dron nazorati", "Analitika"]
)
add_textbox(slide, Inches(4.95), Inches(2.0), Inches(2.2), Inches(2.8),
            "Diagramma mazmuni:\n"
            "• BIM eng katta ulushni egallaydi, chunki u barcha bo‘limlarni bog‘laydi.\n"
            "• Sensorli kuzatuv operativ boshqaruvni ta’minlaydi.\n"
            "• Robotlashtirish ko‘proq takroriy va xavfli vazifalarda samarali.\n"
            "• Dron va tahlil vositalari nazoratning tezkorligini oshiradi.",
            11.8, DARK, False, BODY_FONT)
add_rect(slide, Inches(8.0), Inches(1.8), Inches(4.7), Inches(4.6), SOFT, LIGHT_GRAY, 1.0, radius=True)
add_textbox(slide, Inches(8.25), Inches(2.05), Inches(4.15), Inches(0.4), "Amaliy izoh", 18, DARK, True, TITLE_FONT)
add_textbox(slide, Inches(8.25), Inches(2.55), Inches(3.95), Inches(3.2),
            "Qurilish tashkilotlari odatda avtomatlashtirishni BIM joriy etishdan boshlaydi. Keyin sensorlar va mobil nazorat vositalari qo‘shiladi. "
            "Robotlashtirish esa kapital qiymati yuqoriroq bo‘lgani sababli aniq iqtisodiy hisob-kitob asosida bosqichma-bosqich joriy qilinadi. "
            "Shuning uchun texnologiyalar ulushi bir xil emas; ular korxonaning tayyorgarlik darajasi, ob’ekt turi va byudjetiga bog‘liq ravishda shakllanadi.",
            12.2, DARK, False, BODY_FONT)
add_footer(slide, 5)

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Afzalliklar va kutiladigan natijalar", "Ishlab chiqarish samaradorligi, sifat va boshqaruvdagi ijobiy ta’sirlar", "topline")
add_rect(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(4.95), WHITE, LIGHT_GRAY, 1.0, radius=True)

benefits = [
    ("Tezlik", "Ishlar ketma-ketligi aniq boshqarilib, kutish va takroriy bajarish holatlari kamayadi."),
    ("Sifat", "Sensor, lazer va standart chek-listlar xatolikni erta bosqichda aniqlashga yordam beradi."),
    ("Xavfsizlik", "Xavfli hududlarda odam o‘rniga texnika yoki masofaviy monitoring ishlatiladi."),
    ("Shaffoflik", "Rahbariyat uchun progress, xarajat va bandlik bo‘yicha ishonchli ma’lumot paydo bo‘ladi."),
    ("Tejamkorlik", "Material isrofi, yoqilg‘i sarfi va bekor turish vaqtlarini kamaytirish imkoni yaratiladi."),
    ("Barqarorlik", "Energiya va resurslardan oqilona foydalanish ekologik ko‘rsatkichlarni yaxshilaydi."),
]
for i, (t, txt) in enumerate(benefits):
    col = i % 3
    row = i // 3
    x = Inches(0.9) + col * Inches(3.95)
    y = Inches(2.0) + row * Inches(2.15)
    fillc = GREEN_LIGHT if col == 0 else BLUE_LIGHT if col == 1 else SOFT
    linec = GREEN if col == 0 else BLUE if col == 1 else MID_GRAY
    add_rect(slide, x, y, Inches(3.55), Inches(1.75), fillc, linec, 1.0, radius=True)
    add_textbox(slide, x + Inches(0.14), y + Inches(0.14), Inches(3.1), Inches(0.32), t, 15, DARK, True, TITLE_FONT)
    add_textbox(slide, x + Inches(0.14), y + Inches(0.52), Inches(3.12), Inches(0.92), txt, 11.5, DARK, False, BODY_FONT)

add_footer(slide, 6)

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Qiyinchiliklar va xavflar", "Joriy etish jarayonida e’tibor talab qiladigan masalalar", "left")
add_rect(slide, Inches(0.7), Inches(1.8), Inches(3.8), Inches(4.8), GREEN_LIGHT, GREEN, 1.0, radius=True)
add_textbox(slide, Inches(0.95), Inches(2.05), Inches(3.25), Inches(0.38), "Tashkiliy muammolar", 17, DARK, True, TITLE_FONT)
add_paragraphs_box(
    slide, Inches(0.9), Inches(2.5), Inches(3.2), Inches(3.6),
    [
        "Xodimlarning yangi tizimga moslashuvi vaqt talab qiladi.",
        "Bo‘limlar o‘rtasida ma’lumot almashish standarti bo‘lmasa, platformalar bo‘linib qoladi.",
        "Rahbariyat tomonidan qat’iy KPI va ichki tartib bo‘lmasa, raqamli tizim formallikka aylanadi."
    ],
    font_size=12, color=DARK, bullet=True, bg=None, line=None, radius=False
)
add_rect(slide, Inches(4.8), Inches(1.8), Inches(3.8), Inches(4.8), BLUE_LIGHT, BLUE, 1.0, radius=True)
add_textbox(slide, Inches(5.05), Inches(2.05), Inches(3.25), Inches(0.38), "Texnik cheklovlar", 17, DARK, True, TITLE_FONT)
add_paragraphs_box(
    slide, Inches(5.0), Inches(2.5), Inches(3.2), Inches(3.6),
    [
        "Uskunalar, sensorlar va dasturiy ta’minot integratsiyasi murakkab bo‘lishi mumkin.",
        "Internet yoki lokal tarmoq sifati sust bo‘lsa, real vaqt monitoring samarasi pasayadi.",
        "Texnik xizmat ko‘rsatish va kalibrovka rejimi bo‘lmasa, ma’lumotlar ishonchliligi yomonlashadi."
    ],
    font_size=12, color=DARK, bullet=True, bg=None, line=None, radius=False
)
add_rect(slide, Inches(8.9), Inches(1.8), Inches(3.8), Inches(4.8), SOFT, MID_GRAY, 1.0, radius=True)
add_textbox(slide, Inches(9.15), Inches(2.05), Inches(3.25), Inches(0.38), "Moliyaviy va huquqiy risklar", 17, DARK, True, TITLE_FONT)
add_paragraphs_box(
    slide, Inches(9.1), Inches(2.5), Inches(3.2), Inches(3.6),
    [
        "Boshlang‘ich investitsiya yuqori ko‘rinishi mumkin, ayniqsa kichik loyihalarda.",
        "Texnologiya tanlashda noto‘g‘ri vendor yoki yopiq platforma kelajakdagi moslashuvchanlikni cheklaydi.",
        "Raqamli hujjat aylanishi va ma’lumot xavfsizligi bo‘yicha siyosat ishlab chiqilishi zarur."
    ],
    font_size=12, color=DARK, bullet=True, bg=None, line=None, radius=False
)
add_footer(slide, 7)

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Amaliy qo‘llash senariylari", "Turli bino turlarida avtomatlashtirish qanday ishlaydi", "band")
scenarios = [
    ("Ko‘p qavatli turar joy", "BIM asosida konstruktiv va muhandislik tizimlari oldindan muvofiqlashtiriladi; beton quyish, material oqimi va qavatlar kesimida progress nazorati avtomatlashtiriladi."),
    ("Sanoat binolari", "Katta hajmli metall konstruksiyalar, og‘ir uskunalar va logistika yo‘llari uchun sensorli xavfsizlik hamda avtonom tashish yechimlari qo‘llanadi."),
    ("Savdo va ofis markazlari", "MEP tizimlari murakkab bo‘lgani sababli kolliziya nazorati, montaj ketma-ketligi va sifat hujjatlashtiruvi raqamli boshqariladi."),
    ("Ijtimoiy infratuzilma", "Maktab, shifoxona va davlat obyektlarida qurilish sifati, muddat va shaffof hisobot mexanizmlari ustuvor bo‘lib, avtomatlashtirish aynan shu ko‘rsatkichlarni yaxshilaydi.")
]
for i, (t, txt) in enumerate(scenarios):
    x = Inches(0.8) if i % 2 == 0 else Inches(6.75)
    y = Inches(1.7) + (i // 2) * Inches(2.35)
    fillc = GREEN_LIGHT if i in (0, 3) else BLUE_LIGHT if i == 1 else SOFT
    linec = GREEN if i in (0, 3) else BLUE if i == 1 else MID_GRAY
    add_rect(slide, x, y, Inches(5.75), Inches(1.9), fillc, linec, 1.0, radius=True)
    add_badge(slide, x + Inches(0.16), y + Inches(0.16), Inches(0.55), Inches(0.32), str(i+1), WHITE, linec, linec)
    add_textbox(slide, x + Inches(0.82), y + Inches(0.14), Inches(4.55), Inches(0.3), t, 15, DARK, True, TITLE_FONT)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.58), Inches(5.2), Inches(1.05), txt, 11.5, DARK, False, BODY_FONT)
add_footer(slide, 8)

# Slide 9
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Joriy etish bo‘yicha tavsiyalar", "Bosqichma-bosqich yondashuv orqali xavfni kamaytirish", "topline")
left_x = Inches(0.75)
right_x = Inches(6.8)
items_left = [
    ("1-bosqich: audit", "Mavjud jarayonlar, uskunalar, kadrlar va ma’lumot oqimlari bo‘yicha real holat tahlil qilinadi."),
    ("2-bosqich: pilot loyiha", "Barcha obyektlarda emas, avval bitta nazoratli maydonda eng foydali texnologiya sinovdan o‘tkaziladi."),
    ("3-bosqich: KPI", "Muddat, sifat, xavfsizlik, material isrofi va qayta ishlash hajmi bo‘yicha o‘lchanadigan mezonlar belgilanadi.")
]
items_right = [
    ("4-bosqich: integratsiya", "BIM, mobil nazorat, sensor va hisobot platformalari yagona boshqaruv mantiqiga keltiriladi."),
    ("5-bosqich: o‘qitish", "Muhandislar, ustalar va operatorlar uchun qisqa, amaliy va natijaga yo‘naltirilgan trening tashkil etiladi."),
    ("6-bosqich: masshtablash", "Pilot loyiha natijalari tasdiqlangach, texnologiyalar boshqa obyektlarga bosqichma-bosqich ko‘chiriladi.")
]
for i, (t, txt) in enumerate(items_left):
    y = Inches(1.75) + i * Inches(1.62)
    add_rect(slide, left_x, y, Inches(5.4), Inches(1.25), GREEN_LIGHT if i != 1 else SOFT, GREEN if i != 1 else MID_GRAY, 1.0, radius=True)
    add_textbox(slide, left_x + Inches(0.18), y + Inches(0.12), Inches(4.8), Inches(0.28), t, 14.5, DARK, True, TITLE_FONT)
    add_textbox(slide, left_x + Inches(0.18), y + Inches(0.48), Inches(4.95), Inches(0.54), txt, 11.2, DARK, False, BODY_FONT)
for i, (t, txt) in enumerate(items_right):
    y = Inches(1.75) + i * Inches(1.62)
    add_rect(slide, right_x, y, Inches(5.4), Inches(1.25), BLUE_LIGHT if i != 2 else SOFT, BLUE if i != 2 else MID_GRAY, 1.0, radius=True)
    add_textbox(slide, right_x + Inches(0.18), y + Inches(0.12), Inches(4.8), Inches(0.28), t, 14.5, DARK, True, TITLE_FONT)
    add_textbox(slide, right_x + Inches(0.18), y + Inches(0.48), Inches(4.95), Inches(0.54), txt, 11.2, DARK, False, BODY_FONT)
add_badge(slide, Inches(5.2), Inches(6.45), Inches(2.9), Inches(0.34), "Pilotdan masshtablashgacha", GREEN_LIGHT, GREEN, GREEN_DARK)
add_footer(slide, 9)

# Slide 10
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_rect(slide, Inches(0.6), Inches(0.7), Inches(12.1), Inches(5.95), WHITE, LIGHT_GRAY, 1.0, radius=True)
add_rect(slide, Inches(0.9), Inches(1.0), Inches(0.18), Inches(0.95), GREEN, GREEN, 0.2, radius=True)
add_textbox(slide, Inches(1.2), Inches(1.0), Inches(5.8), Inches(0.55), "Xulosa", 29, DARK, True, TITLE_FONT)
add_textbox(slide, Inches(1.22), Inches(1.62), Inches(10.8), Inches(1.35),
            "Avtomatlashtirilgan zamonaviy bino qurilish tehnikalari qurilish sohasini tezkor, aniq va shaffof boshqariladigan tizimga aylantirmoqda. "
            "BIM, sensorlar, robotlashtirilgan uskunalar va analitik boshqaruv birgalikda ishlaganda loyiha sifati oshadi, xarajatlar nazorat qilinadi va xavfsizlik darajasi mustahkamlanadi.",
            15, DARK, False, BODY_FONT)
add_rect(slide, Inches(1.15), Inches(3.15), Inches(3.35), Inches(2.3), GREEN_LIGHT, GREEN, 1.0, radius=True)
add_textbox(slide, Inches(1.35), Inches(3.38), Inches(2.85), Inches(0.32), "Kalit g‘oya 1", 16, DARK, True, TITLE_FONT)
add_textbox(slide, Inches(1.35), Inches(3.8), Inches(2.8), Inches(1.25),
            "Texnologiya alohida vosita emas; u loyiha, jadval, sifat va resurslarni birlashtiruvchi boshqaruv falsafasidir.",
            11.5, DARK, False, BODY_FONT)
add_rect(slide, Inches(4.95), Inches(3.15), Inches(3.35), Inches(2.3), BLUE_LIGHT, BLUE, 1.0, radius=True)
add_textbox(slide, Inches(5.15), Inches(3.38), Inches(2.85), Inches(0.32), "Kalit g‘oya 2", 16, DARK, True, TITLE_FONT)
add_textbox(slide, Inches(5.15), Inches(3.8), Inches(2.8), Inches(1.25),
            "Eng to‘g‘ri yo‘l — kichik pilot loyiha, o‘lchab boriladigan KPI va xodimlarni amaliy o‘qitish orqali joriy etishdir.",
            11.5, DARK, False, BODY_FONT)
add_rect(slide, Inches(8.75), Inches(3.15), Inches(2.95), Inches(2.3), SOFT, MID_GRAY, 1.0, radius=True)
add_textbox(slide, Inches(8.95), Inches(3.38), Inches(2.45), Inches(0.32), "Kalit g‘oya 3", 16, DARK, True, TITLE_FONT)
add_textbox(slide, Inches(8.95), Inches(3.8), Inches(2.35), Inches(1.25),
            "Barqaror natija integratsiya, intizom va real vaqt ma’lumotlaridan oqilona foydalanishga bog‘liq.",
            11.5, DARK, False, BODY_FONT)
add_badge(slide, Inches(1.2), Inches(5.92), Inches(2.5), Inches(0.34), f"Taqdimot egasi: {TAQDIMOT_EGASI}", GREEN_LIGHT, GREEN, GREEN_DARK)
add_textbox(slide, Inches(9.0), Inches(5.9), Inches(2.8), Inches(0.32), "E’tiboringiz uchun rahmat", 15, DARK, True, TITLE_FONT)
add_footer(slide, 10)

while len(prs.slides) > SLIDE_COUNT:
    rId = prs.slides._sldIdLst[-1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[-1]

prs.save(OUTPUT_PATH)